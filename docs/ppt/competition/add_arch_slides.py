#!/usr/bin/env python3
"""Append 2 detailed technical architecture slides to competition.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x0F, 0x0F, 0x23)
DARK_BG2 = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
PURPLE = RGBColor(0xA7, 0x8B, 0xF0)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
LAYER_BG = RGBColor(0xF1, 0xF5, 0xF9)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
SUBTLE_TEXT = RGBColor(0xAA, 0xAA, 0xBB)
DB_COLOR = RGBColor(0x05, 0x96, 0x69)
RAG_COLOR = RGBColor(0x6D, 0x28, 0xD9)

# Load existing presentation
prs = Presentation(r'C:\Users\Administrator\Desktop\Project\claudeproject\solution-implementations\ppt\competition\competition.pptx')

def set_bg(slide, color):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill=None, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=Pt(12), color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = sz; p.font.color.rgb = color; p.font.bold = bold; p.font.name = font; p.alignment = align
    return box

def chrome(s, left, right, dark=True):
    c = WHITE if dark else GRAY_TEXT
    tb(s, Inches(0.6), Inches(0.25), Inches(6), Inches(0.3), left, Pt(10), c, font='Consolas')
    tb(s, Inches(10), Inches(0.25), Inches(2.8), Inches(0.3), right, Pt(10), c, align=PP_ALIGN.RIGHT, font='Consolas')
    rect(s, Inches(0.6), Inches(0.6), Inches(12.1), Pt(1), fill=RGBColor(0x33,0x33,0x50) if dark else CARD_BORDER)

def footer(s, left, right):
    tb(s, Inches(0.6), Inches(7.0), Inches(6), Inches(0.35), left, Pt(9), GRAY_TEXT, font='Consolas')
    tb(s, Inches(10), Inches(7.0), Inches(2.8), Inches(0.35), right, Pt(9), GRAY_TEXT, align=PP_ALIGN.RIGHT, font='Consolas')

def layer_label(s, l, t, text, color=BLUE):
    rrect(s, l, t, Inches(1.3), Inches(0.22), fill=color)
    tb(s, l+Inches(0.06), t+Pt(1), Inches(1.2), Inches(0.2), text, Pt(8), WHITE, bold=True, font='Consolas')


# ================================================================
# SLIDE 29: InsightPro 详细技术架构
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, LIGHT_BG)
chrome(s, 'InsightPro · 详细技术架构', 'Act II / Arch Detail · 29 / 30', dark=False)
tb(s, Inches(0.8), Inches(0.8), Inches(8), Inches(0.5), 'InsightPro 全栈技术架构', Pt(32), DARK_TEXT, bold=True)
tb(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.3), 'DeepSearcher RAG + 双AI引擎 + 31个API端点 + 15个页面路由 + 9张SQLite表', Pt(12), GRAY_TEXT)

# === FRONTEND LAYER ===
tb(s, Inches(0.6), Inches(1.75), Inches(5), Inches(0.25), 'Frontend Layer', Pt(9), BLUE, bold=True, font='Consolas')
rect(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.65), fill=LAYER_BG, border=CARD_BORDER)
boxes_fe = [
    ('Next.js 16.2.6', 'App Router 15路由', BLUE),
    ('React 19.2.4', 'TypeScript 5', BLUE),
    ('Tailwind CSS 4', '响应式布局', GREEN),
    ('Recharts 3.8', 'Area/Bar/Line', ORANGE),
    ('ChatAssistant', 'SSE流式对话', PURPLE),
    ('PageTracker', 'sendBeacon 埋点', GRAY_TEXT),
]
for i, (t1, t2, c) in enumerate(boxes_fe):
    x = Inches(0.65 + i * 2.02)
    rrect(s, x, Inches(2.05), Inches(1.95), Inches(0.55), fill=WHITE, border=c)
    tb(s, x+Inches(0.08), Inches(2.08), Inches(1.8), Inches(0.22), t1, Pt(9), DARK_TEXT, bold=True)
    tb(s, x+Inches(0.08), Inches(2.32), Inches(1.8), Inches(0.2), t2, Pt(7.5), GRAY_TEXT)

# Arrow down
tb(s, Inches(6.2), Inches(2.7), Inches(1), Inches(0.25), '▼', Pt(14), GRAY_TEXT, align=PP_ALIGN.CENTER)

# === API LAYER ===
tb(s, Inches(0.6), Inches(3.0), Inches(5), Inches(0.25), 'API Layer (FastAPI 0.136 · Uvicorn 0.48)', Pt(9), GREEN, bold=True, font='Consolas')
rect(s, Inches(0.6), Inches(3.25), Inches(12.1), Inches(0.65), fill=LAYER_BG, border=CARD_BORDER)
boxes_api = [
    ('/api/chat/stream', 'SSE流式客服', PURPLE),
    ('/api/baidu-hotsearch', '百度热搜爬虫', ORANGE),
    ('/api/github-trending', 'GitHub Trending', ORANGE),
    ('/api/market/*', '市场情报RAG', RAG_COLOR),
    ('/api/bidding/*', '招标信息分析', GREEN),
    ('/api/demand/*', '需求挖掘', CYAN),
]
for i, (t1, t2, c) in enumerate(boxes_api):
    x = Inches(0.65 + i * 2.02)
    rrect(s, x, Inches(3.3), Inches(1.95), Inches(0.55), fill=WHITE, border=c)
    tb(s, x+Inches(0.08), Inches(3.33), Inches(1.8), Inches(0.22), t1, Pt(8), DARK_TEXT, bold=True, font='Consolas')
    tb(s, x+Inches(0.08), Inches(3.57), Inches(1.8), Inches(0.2), t2, Pt(7.5), GRAY_TEXT)

# Arrow down
tb(s, Inches(6.2), Inches(3.95), Inches(1), Inches(0.25), '▼', Pt(14), GRAY_TEXT, align=PP_ALIGN.CENTER)

# === AI ENGINE LAYER ===
tb(s, Inches(0.6), Inches(4.25), Inches(5), Inches(0.25), 'AI Engine Layer (双AI引擎 + RAG增强)', Pt(9), RAG_COLOR, bold=True, font='Consolas')
rect(s, Inches(0.6), Inches(4.5), Inches(5.85), Inches(1.55), fill=LAYER_BG, border=RAG_COLOR)
layer_label(s, Inches(0.68), Inches(4.55), 'DeepSeek API', RAG_COLOR)
tb(s, Inches(0.68), Inches(4.85), Inches(5.5), Inches(0.2), 'api.deepseek.com · deepseek-chat · OpenAI SDK兼容', Pt(8), DARK_TEXT, font='Consolas')
tb(s, Inches(0.68), Inches(5.05), Inches(5.5), Inches(0.2), '用途：任务分析 / 招标分析 / 市场洞察 / 需求报告 / JSON格式输出', Pt(8), GRAY_TEXT)

rect(s, Inches(0.6), Inches(5.3), Inches(5.85), Inches(0.7), fill=WHITE, border=GREEN)
layer_label(s, Inches(0.68), Inches(5.35), 'ModelArts MaaS', GREEN)
tb(s, Inches(0.68), Inches(5.6), Inches(5.5), Inches(0.35), 'api.modelarts-maas.com · deepseek-v3.2 · httpx直连 · 流式SSE + 非流式\n用途：站内AI智能客服问答（高频）', Pt(8), DARK_TEXT)

# RAG block (right side)
rect(s, Inches(6.65), Inches(4.5), Inches(6.05), Inches(1.55), fill=LAYER_BG, border=RAG_COLOR)
layer_label(s, Inches(6.73), Inches(4.55), 'DeepSearcher RAG', RAG_COLOR)
tb(s, Inches(6.73), Inches(4.82), Inches(5.8), Inches(0.2), 'Embedding: BAAI/bge-m3 · 1024维 · 中英双语 · 本地SentenceTransformer', Pt(8), DARK_TEXT, font='Consolas')
tb(s, Inches(6.73), Inches(5.02), Inches(5.8), Inches(0.2), 'VectorDB: Qdrant 本地文件持久化 (deep_searcher_qdrant/)', Pt(8), DARK_TEXT, font='Consolas')
tb(s, Inches(6.73), Inches(5.22), Inches(5.8), Inches(0.2), '5个向量集合：', Pt(8), RAG_COLOR, bold=True)
collections = ['platform_knowledge', 'demand_signals', 'competitive_news', 'github_trending', 'policies']
for i, coll in enumerate(collections):
    x = Inches(6.73 + (i % 3) * 2.0)
    y = Inches(5.45) + Inches(0.2 if i >= 3 else 0)
    rrect(s, x, y, Inches(1.85), Inches(0.2), fill=WHITE, border=RAG_COLOR)
    tb(s, x+Inches(0.05), y+Pt(1), Inches(1.75), Inches(0.18), coll, Pt(7), RAG_COLOR, font='Consolas')

# Arrow down
tb(s, Inches(6.2), Inches(6.1), Inches(1), Inches(0.25), '▼', Pt(14), GRAY_TEXT, align=PP_ALIGN.CENTER)

# === DATA LAYER ===
tb(s, Inches(0.6), Inches(6.4), Inches(5), Inches(0.25), 'Data Layer (SQLite 9表 + Supabase预留)', Pt(9), DB_COLOR, bold=True, font='Consolas')
rect(s, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.8), fill=LAYER_BG, border=DB_COLOR)
layer_label(s, Inches(0.68), Inches(6.7), 'SQLite DB', DB_COLOR)

tables = ['github_trending', 'baidu_hotsearch', 'scrape_log', 'page_visits', 'email_subscribers',
          'bidding_opportunities', 'competitor_news', 'demand_signals', 'demand_reports']
for i, tbl in enumerate(tables):
    col = i % 5
    row = i // 5
    x = Inches(0.68 + col * 2.42)
    y = Inches(6.95 + row * 0.22)
    rrect(s, x, y, Inches(2.3), Inches(0.2), fill=WHITE, border=DB_COLOR)
    tb(s, x+Inches(0.05), y+Pt(1), Inches(2.2), Inches(0.18), tbl, Pt(7), DB_COLOR, font='Consolas')

# Supabase
layer_label(s, Inches(7.5), Inches(6.7), 'Supabase PG', BLUE)
tb(s, Inches(7.5), Inches(6.95), Inches(5), Inches(0.2), 'Prisma 7.8 Schema: User / InsightTask / DataSource / Report / DailyInsight', Pt(8), DARK_TEXT, font='Consolas')
tb(s, Inches(7.5), Inches(7.15), Inches(5), Inches(0.2), 'PgBouncer :6543 (Transaction) + Direct :5432 (Session) · 已配置·预留扩展', Pt(7.5), GRAY_TEXT, font='Consolas')

# === SCHEDULER & CRAWLER STRATEGY ===
rect(s, Inches(0.6), Inches(7.5), Inches(12.1), Pt(1), fill=CARD_BORDER)
tb(s, Inches(0.8), Inches(7.55), Inches(12), Inches(0.2), 'APScheduler: 09:00 GitHub | 08:30 招标 | 08:00 需求 | 09:05 邮件 | 03:00 清理  ·  降级策略: 实时爬取→SQLite降级  ·  双AI: DeepSeek(低频分析) + ModelArts(高频客服)', Pt(8), GRAY_TEXT, font='Consolas')

footer(s, 'InsightPro · 详细技术架构', 'Architecture Detail')

print("  Slide 29: InsightPro Technical Architecture")


# ================================================================
# SLIDE 30: SAC 解决方案即代码 详细技术架构
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
chrome(s, 'SAC · 解决方案即代码 详细架构', 'Act III / Arch Detail · 30 / 30')
tb(s, Inches(0.8), Inches(0.8), Inches(10), Inches(0.5), 'SAC 全栈技术架构', Pt(32), WHITE, bold=True)
tb(s, Inches(0.8), Inches(1.3), Inches(11), Inches(0.3), '4个AI Skills · 9个Python脚本 · 7个解决方案 · CN/HK双区域 · RFS一键部署', Pt(12), SUBTLE_TEXT)

# === SOLUTIONS LAYER ===
tb(s, Inches(0.6), Inches(1.75), Inches(5), Inches(0.25), 'Solutions Layer (7个方案实践)', Pt(9), GOLD, bold=True, font='Consolas')
rect(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.8), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))

sols = [
    ('Headroom OpenCode', 'HK · pip直装', BLUE),
    ('Headroom Claude Code', 'HK · pip直装', BLUE),
    ('LiteLLM', 'CN/HK · Docker', GREEN),
    ('Supabase', 'CN · Docker', GREEN),
    ('OpenHands', 'CN · Docker', PURPLE),
    ('CodeWhale', 'CN · Docker', PURPLE),
    ('AiToEarn', 'CN/HK · Docker', ORANGE),
]
for i, (name, mode, c) in enumerate(sols):
    x = Inches(0.65 + i * 1.72)
    rrect(s, x, Inches(2.05), Inches(1.65), Inches(0.7), fill=RGBColor(0x25,0x25,0x40), border=c)
    tb(s, x+Inches(0.08), Inches(2.08), Inches(1.5), Inches(0.22), name, Pt(9), WHITE, bold=True)
    tb(s, x+Inches(0.08), Inches(2.35), Inches(1.5), Inches(0.25), mode, Pt(7), SUBTLE_TEXT)

# Arrow down
tb(s, Inches(6.2), Inches(2.85), Inches(1), Inches(0.25), '▼', Pt(14), SUBTLE_TEXT, align=PP_ALIGN.CENTER)

# === TERRAFORM LAYER ===
tb(s, Inches(0.6), Inches(3.1), Inches(8), Inches(0.25), 'Terraform Layer (HCL模板 · 7核心变量 · 唯一命名)', Pt(9), GREEN, bold=True, font='Consolas')
rect(s, Inches(0.6), Inches(3.35), Inches(7.5), Inches(1.3), fill=DARK_BG2, border=GREEN)
layer_label(s, Inches(0.68), Inches(3.4), 'IaC Template', GREEN)

tf_res = [
    ('huaweicloud_vpc', '172.16.0.0/16', BLUE),
    ('huaweicloud_vpc_subnet', '172.16.1.0/24', BLUE),
    ('networking_secgroup', 'SSH+Proxy+ICMP', GREEN),
    ('vpc_eip', '300Mbit/s 流量', ORANGE),
    ('compute_instance', 'Flexus X · Ubuntu', PURPLE),
]
for i, (res, desc, c) in enumerate(tf_res):
    col = i % 3
    row = i // 3
    x = Inches(0.68 + col * 2.5)
    y = Inches(3.72 + row * 0.4)
    rrect(s, x, y, Inches(2.35), Inches(0.35), fill=RGBColor(0x25,0x25,0x40), border=c)
    tb(s, x+Inches(0.05), y+Pt(1), Inches(2.25), Inches(0.15), res, Pt(7.5), c, font='Consolas')
    tb(s, x+Inches(0.05), y+Pt(13), Inches(2.25), Inches(0.15), desc, Pt(7), SUBTLE_TEXT)

# user_data flow
rect(s, Inches(8.3), Inches(3.35), Inches(4.4), Inches(1.3), fill=DARK_BG2, border=GOLD)
layer_label(s, Inches(8.38), Inches(3.4), 'user_data Flow', GOLD)
steps_ud = ['① chpasswd 重置密码', '② wget OBS 下载 install.sh', '③ bash 执行安装脚本', '④ rm -rf 清理临时文件']
for i, step in enumerate(steps_ud):
    tb(s, Inches(8.38), Inches(3.75 + i * 0.22), Inches(4.2), Inches(0.2), step, Pt(8), WHITE)
tb(s, Inches(8.38), Inches(4.65), Inches(4.2), Inches(0.2), '模板-脚本分离原则 · user_data最小化', Pt(7), GOLD, font='Consolas')

# Arrow down
tb(s, Inches(6.2), Inches(4.7), Inches(1), Inches(0.25), '▼', Pt(14), SUBTLE_TEXT, align=PP_ALIGN.CENTER)

# === INSTALL SCRIPT LAYER ===
tb(s, Inches(0.6), Inches(5.0), Inches(8), Inches(0.25), 'Install Script Layer (5阶段 · 4模式)', Pt(9), ORANGE, bold=True, font='Consolas')
rect(s, Inches(0.6), Inches(5.25), Inches(12.1), Inches(0.75), fill=DARK_BG2, border=ORANGE)

stages = [
    ('Stage 1', '系统准备', 'dpkg\napt update\nca-certs', GOLD),
    ('Stage 2', 'Docker', 'CN:华为云镜像\nHK:官方源\ncompose up', BLUE),
    ('Stage 3', '应用配置', '目录权限\n配置生成\ncrontab备份', GREEN),
    ('Stage 4', '启动验证', 'docker up\n健康检查\n120s重试', PURPLE),
    ('Stage 5', 'OpenCode', 'Node.js 18+\nnpm opencode\nheadroom-ai', CYAN),
]
for i, (stage, name, desc, c) in enumerate(stages):
    x = Inches(0.65 + i * 2.44)
    rrect(s, x, Inches(5.3), Inches(2.35), Inches(0.65), fill=RGBColor(0x25,0x25,0x40), border=c)
    tb(s, x+Inches(0.05), Inches(5.32), Inches(2.25), Inches(0.15), f'{stage} · {name}', Pt(8), c, bold=True)
    tb(s, x+Inches(0.05), Inches(5.5), Inches(2.25), Inches(0.35), desc, Pt(7), SUBTLE_TEXT)

# Arrow down
tb(s, Inches(6.2), Inches(6.05), Inches(1), Inches(0.25), '▼', Pt(14), SUBTLE_TEXT, align=PP_ALIGN.CENTER)

# === ECS RUNTIME LAYER ===
tb(s, Inches(0.6), Inches(6.3), Inches(8), Inches(0.25), 'ECS Runtime (Flexus X · Ubuntu 24.04 · x1.2u.4g)', Pt(9), BLUE, bold=True, font='Consolas')
rect(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.7), fill=DARK_BG2, border=BLUE)

rt_items = [
    ('OpenCode CLI', 'npm -g\nopencode-ai@latest\nterminal-native', BLUE),
    ('Headroom Proxy', 'pip install\nheadroom-ai\nport :8787', GREEN),
    ('Context Engine', 'AST感知压缩\n60-95% Token节省\n7语言支持', RAG_COLOR),
    ('MaaS Endpoint', 'modelarts-maas.com\ndeepseek-v3.2\nclaude切换', GOLD),
    ('Docker Services', 'LiteLLM:4000\nPG:5432\nProm:9090', ORANGE),
    ('Monitoring', '/metrics\n/stats\nPrometheus兼容', CYAN),
]
for i, (name, desc, c) in enumerate(rt_items):
    x = Inches(0.65 + i * 2.02)
    rrect(s, x, Inches(6.6), Inches(1.95), Inches(0.6), fill=RGBColor(0x25,0x25,0x40), border=c)
    tb(s, x+Inches(0.05), Inches(6.62), Inches(1.85), Inches(0.15), name, Pt(8), c, bold=True)
    tb(s, x+Inches(0.05), Inches(6.8), Inches(1.85), Inches(0.3), desc, Pt(7), SUBTLE_TEXT)

# === AI SKILLS (right side) ===
rect(s, Inches(9.5), Inches(1.75), Inches(3.2), Inches(2.8), fill=DARK_BG2, border=GOLD)
tb(s, Inches(9.6), Inches(1.8), Inches(3), Inches(0.25), 'AI Skills 技能体系', Pt(10), GOLD, bold=True)
skills = [
    ('rfs-practices', '创建/修复RFS模板', '10决策点 · 27 Pitfall'),
    ('solution-extractor', '方案页面→结构化Excel', 'browser→snapshot→xlsx'),
    ('page-enhance', '方案页面文案增强', '10维度优化检查'),
    ('deep-search-synth', '3层深度调研', '多源交叉验证'),
]
for i, (name, desc, detail) in enumerate(skills):
    y = Inches(2.1 + i * 0.6)
    tb(s, Inches(9.6), y, Inches(3), Inches(0.2), name, Pt(9), WHITE, bold=True, font='Consolas')
    tb(s, Inches(9.6), y+Inches(0.18), Inches(3), Inches(0.15), desc, Pt(8), SUBTLE_TEXT)
    tb(s, Inches(9.6), y+Inches(0.33), Inches(3), Inches(0.15), detail, Pt(7), GRAY_TEXT, font='Consolas')

# === SCRIPTS (right side) ===
rect(s, Inches(9.5), Inches(4.7), Inches(3.2), Inches(2.5), fill=DARK_BG2, border=GREEN)
tb(s, Inches(9.6), Inches(4.75), Inches(3), Inches(0.25), 'Python/Bash 工具链', Pt(10), GREEN, bold=True)
scripts = [
    ('validate_template.py', 'RFS模板结构验证'),
    ('gen_docx.py', 'README→Word中英文'),
    ('gen_sac_docx.py', 'SAC方案报告生成'),
    ('generate_extension.py', '.extension配置生成'),
    ('package_solution.sh', '方案打包归档(.zip)'),
]
for i, (name, desc) in enumerate(scripts):
    y = Inches(5.1 + i * 0.4)
    tb(s, Inches(9.6), y, Inches(3), Inches(0.18), name, Pt(8), WHITE, bold=True, font='Consolas')
    tb(s, Inches(9.6), y+Inches(0.17), Inches(3), Inches(0.15), desc, Pt(7.5), SUBTLE_TEXT)

footer(s, 'SAC · 详细技术架构', 'Solution as Code')

# Save
prs.save(r'C:\Users\Administrator\Desktop\Project\claudeproject\solution-implementations\ppt\competition\competition.pptx')
print(f"Saved! Total slides: {len(prs.slides)}")
