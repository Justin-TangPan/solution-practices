#!/usr/bin/env python3
"""Generate 2 detailed technical architecture slides (standalone)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
DARK_BG = RGBColor(0x0F, 0x0F, 0x23)
DARK_BG2 = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
PURPLE = RGBColor(0xA7, 0x8B, 0xF0)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
RAG = RGBColor(0x6D, 0x28, 0xD9)
DB = RGBColor(0x05, 0x96, 0x69)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
LAYER_BG = RGBColor(0xF1, 0xF5, 0xF9)
SUBTLE = RGBColor(0xAA, 0xAA, 0xBB)
BOX_DARK = RGBColor(0x25, 0x25, 0x40)
RULE = RGBColor(0x40, 0x40, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──
def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def box(s, l, t, w, h, fill=None, border=None, r=False):
    sh = MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE
    o = s.shapes.add_shape(sh, l, t, w, h)
    if fill: o.fill.solid(); o.fill.fore_color.rgb = fill
    else: o.fill.background()
    if border: o.line.color.rgb = border; o.line.width = Pt(0.75)
    else: o.line.fill.background()
    return o

def txt(s, l, t, w, h, text, sz=Pt(12), c=DARK_TEXT, b=False, a=PP_ALIGN.LEFT, f='Microsoft YaHei'):
    o = s.shapes.add_textbox(l, t, w, h)
    p = o.text_frame.paragraphs[0]; p.text = text
    p.font.size = sz; p.font.color.rgb = c; p.font.bold = b; p.font.name = f; p.alignment = a
    o.text_frame.word_wrap = True
    return o

def header(s, left, right, dark=False):
    fc = WHITE if dark else GRAY
    txt(s, Inches(0.6), Inches(0.25), Inches(6), Inches(0.3), left, Pt(10), fc, f='Consolas')
    txt(s, Inches(10), Inches(0.25), Inches(2.8), Inches(0.3), right, Pt(10), fc, a=PP_ALIGN.RIGHT, f='Consolas')
    box(s, Inches(0.6), Inches(0.6), Inches(12.1), Pt(1), fill=RGBColor(0x33,0x33,0x50) if dark else BORDER)

def foot(s, l, r):
    txt(s, Inches(0.6), Inches(7.0), Inches(6), Inches(0.35), l, Pt(9), GRAY, f='Consolas')
    txt(s, Inches(10), Inches(7.0), Inches(2.8), Inches(0.35), r, Pt(9), GRAY, a=PP_ALIGN.RIGHT, f='Consolas')

def tag(s, l, t, text, color=BLUE):
    box(s, l, t, Inches(1.2), Inches(0.2), r=True, fill=color)
    txt(s, l+Inches(0.05), t, Inches(1.1), Inches(0.2), text, Pt(7.5), WHITE, b=True, f='Consolas')

def small_box(s, l, t, w, h, line1, line2, color=BLUE, dark=False):
    bg_c = BOX_DARK if dark else WHITE
    box(s, l, t, w, h, r=True, fill=bg_c, border=color)
    txt(s, l+Inches(0.06), t+Inches(0.03), w-Inches(0.12), Inches(0.18), line1, Pt(8.5), WHITE if dark else DARK_TEXT, b=True)
    txt(s, l+Inches(0.06), t+Inches(0.22), w-Inches(0.12), Inches(0.16), line2, Pt(7), SUBTLE if dark else GRAY)

def layer_title(s, l, t, text, color=BLUE):
    txt(s, l, t, Inches(6), Inches(0.22), text, Pt(9), color, b=True, f='Consolas')

def v_arrow(s, x, y):
    txt(s, x, y, Inches(0.4), Inches(0.2), '▼', Pt(11), GRAY, a=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: InsightPro
# ══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s1, LIGHT_BG)
header(s1, 'InsightPro · 详细技术架构', '29 / 30', dark=False)

txt(s1, Inches(0.8), Inches(0.75), Inches(10), Inches(0.45), 'InsightPro 全栈技术架构', Pt(30), DARK_TEXT, b=True)
txt(s1, Inches(0.8), Inches(1.15), Inches(11), Inches(0.25),
    'DeepSearcher RAG + 双AI引擎 + 31个API端点 + 15个页面路由 + 9张SQLite表 + APScheduler',
    Pt(11), GRAY)

# ── Frontend ──
layer_title(s1, Inches(0.6), Inches(1.55), 'Frontend Layer', BLUE)
box(s1, Inches(0.6), Inches(1.78), Inches(12.1), Inches(0.55), fill=LAYER_BG, border=BORDER)
for i,(a,b,c) in enumerate([
    ('Next.js 16.2.6','App Router 15路由',BLUE),('React 19.2.4','TypeScript 5',BLUE),
    ('Tailwind CSS 4','响应式布局',GREEN),('Recharts 3.8','Area/Bar/Line',ORANGE),
    ('ChatAssistant','SSE流式对话',PURPLE),('PageTracker','sendBeacon埋点',GRAY)]):
    x = Inches(0.65 + i*2.02)
    small_box(s1, x, Inches(1.82), Inches(1.95), Inches(0.44), a, b, c)

v_arrow(s1, Inches(6.2), Inches(2.36))

# ── API ──
layer_title(s1, Inches(0.6), Inches(2.58), 'API Layer  (FastAPI 0.136 · Uvicorn 0.48 · 31端点)', GREEN)
box(s1, Inches(0.6), Inches(2.8), Inches(12.1), Inches(0.55), fill=LAYER_BG, border=BORDER)
for i,(a,b,c) in enumerate([
    ('/api/chat/stream','SSE流式客服',PURPLE),('/api/baidu-hotsearch','百度热搜爬虫',ORANGE),
    ('/api/github-trending','GitHub Trending',ORANGE),('/api/market/*','市场情报RAG',RAG),
    ('/api/bidding/*','招标分析',GREEN),('/api/demand/*','需求挖掘',CYAN)]):
    x = Inches(0.65 + i*2.02)
    small_box(s1, x, Inches(2.84), Inches(1.95), Inches(0.44), a, b, c)

v_arrow(s1, Inches(6.2), Inches(3.38))

# ── AI Engine ──
layer_title(s1, Inches(0.6), Inches(3.6), 'AI Engine Layer  (双AI引擎 + RAG增强)', RAG)

# DeepSeek
box(s1, Inches(0.6), Inches(3.82), Inches(5.85), Inches(0.6), fill=LAYER_BG, border=RAG)
tag(s1, Inches(0.68), Inches(3.87), 'DeepSeek API', RAG)
txt(s1, Inches(0.68), Inches(4.1), Inches(5.5), Inches(0.14), 'api.deepseek.com · deepseek-chat · OpenAI SDK兼容 · JSON输出', Pt(8), DARK_TEXT, f='Consolas')
txt(s1, Inches(0.68), Inches(4.25), Inches(5.5), Inches(0.14), '任务分析 / 招标分析 / 市场洞察 / 需求报告（低频结构化）', Pt(7.5), GRAY)

# ModelArts
box(s1, Inches(0.6), Inches(4.47), Inches(5.85), Inches(0.55), fill=WHITE, border=GREEN)
tag(s1, Inches(0.68), Inches(4.52), 'ModelArts MaaS', GREEN)
txt(s1, Inches(0.68), Inches(4.75), Inches(5.5), Inches(0.14), 'api.modelarts-maas.com · deepseek-v3.2 · httpx直连 · 流式SSE+非流式', Pt(8), DARK_TEXT, f='Consolas')
txt(s1, Inches(0.68), Inches(4.9), Inches(5.5), Inches(0.14), '站内AI智能客服问答（高频实时）', Pt(7.5), GRAY)

# RAG
box(s1, Inches(6.65), Inches(3.82), Inches(6.05), Inches(1.2), fill=LAYER_BG, border=RAG)
tag(s1, Inches(6.73), Inches(3.87), 'DeepSearcher RAG', RAG)
txt(s1, Inches(6.73), Inches(4.1), Inches(5.8), Inches(0.14), 'Embedding: BAAI/bge-m3 · 1024维 · 中英双语 · 本地SentenceTransformer', Pt(8), DARK_TEXT, f='Consolas')
txt(s1, Inches(6.73), Inches(4.26), Inches(5.8), Inches(0.14), 'VectorDB: Qdrant 本地文件持久化 (deep_searcher_qdrant/)', Pt(8), DARK_TEXT, f='Consolas')
txt(s1, Inches(6.73), Inches(4.42), Inches(3), Inches(0.14), '5个向量集合：', Pt(8), RAG, b=True)
for i, coll in enumerate(['platform_knowledge','demand_signals','competitive_news','github_trending','policies']):
    x = Inches(6.73 + (i%3)*1.95)
    y = Inches(4.58) + (Inches(0.2) if i>=3 else Inches(0))
    box(s1, x, y, Inches(1.8), Inches(0.17), r=True, fill=WHITE, border=RAG)
    txt(s1, x+Inches(0.04), y, Inches(1.7), Inches(0.17), coll, Pt(7), RAG, f='Consolas')

v_arrow(s1, Inches(6.2), Inches(5.1))

# ── Data Layer ──
layer_title(s1, Inches(0.6), Inches(5.32), 'Data Layer  (SQLite 9表 + Supabase PostgreSQL 预留)', DB)

box(s1, Inches(0.6), Inches(5.55), Inches(7.5), Inches(0.7), fill=LAYER_BG, border=DB)
tag(s1, Inches(0.68), Inches(5.6), 'SQLite trending.db', DB)
for i, tbl in enumerate(['github_trending','baidu_hotsearch','scrape_log','page_visits','email_subscribers','bidding_opportunities','competitor_news','demand_signals','demand_reports']):
    x = Inches(0.68 + (i%5)*1.47)
    y = Inches(5.85 + (i//5)*0.18)
    box(s1, x, y, Inches(1.4), Inches(0.16), r=True, fill=WHITE, border=DB)
    txt(s1, x+Inches(0.04), y, Inches(1.3), Inches(0.16), tbl, Pt(7), DB, f='Consolas')

box(s1, Inches(8.3), Inches(5.55), Inches(4.4), Inches(0.7), fill=WHITE, border=BLUE)
tag(s1, Inches(8.38), Inches(5.6), 'Supabase PG', BLUE)
txt(s1, Inches(8.38), Inches(5.82), Inches(4.2), Inches(0.14), 'Prisma 7.8: User / InsightTask / DataSource / Report / DailyInsight', Pt(8), DARK_TEXT, f='Consolas')
txt(s1, Inches(8.38), Inches(5.98), Inches(4.2), Inches(0.14), 'PgBouncer :6543 + Direct :5432 · ap-southeast-1 · 已配置·预留扩展', Pt(7.5), GRAY, f='Consolas')

# Bottom
box(s1, Inches(0.6), Inches(6.35), Inches(12.1), Pt(1), fill=BORDER)
txt(s1, Inches(0.8), Inches(6.4), Inches(12), Inches(0.2),
    'APScheduler: 09:00 GitHub | 08:30 招标 | 08:00 需求 | 09:05 邮件 | 03:00 清理  ·  降级: 实时爬取→SQLite  ·  双AI: DeepSeek(低频) + ModelArts(高频)',
    Pt(8), GRAY, f='Consolas')

foot(s1, 'InsightPro · 详细技术架构', 'Architecture Detail')


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: SAC
# ══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
bg(s2, DARK_BG)
header(s2, 'SAC · 解决方案实践 详细架构', '30 / 30', dark=True)

txt(s2, Inches(0.8), Inches(0.75), Inches(10), Inches(0.45), 'SAC 全栈技术架构', Pt(30), WHITE, b=True)
txt(s2, Inches(0.8), Inches(1.15), Inches(11), Inches(0.25),
    '4个AI Skills · 9个Python脚本 · 7个解决方案 · CN/HK双区域 · RFS一键部署',
    Pt(11), SUBTLE)

# ── Solutions ──
layer_title(s2, Inches(0.6), Inches(1.55), 'Solutions Layer  (7个方案实践)', GOLD)
box(s2, Inches(0.6), Inches(1.78), Inches(12.1), Inches(0.55), fill=DARK_BG2, border=RULE)
for i,(a,b,c) in enumerate([('Headroom OpenCode','HK',BLUE),('Headroom ClaudeCode','HK',BLUE),
    ('LiteLLM','CN/HK',GREEN),('Supabase','CN',GREEN),('OpenHands','CN',PURPLE),
    ('CodeWhale','CN',PURPLE),('AiToEarn','CN/HK',ORANGE)]):
    x = Inches(0.65 + i*1.72)
    box(s2, x, Inches(1.82), Inches(1.65), Inches(0.44), r=True, fill=BOX_DARK, border=c)
    txt(s2, x+Inches(0.08), Inches(1.84), Inches(1.5), Inches(0.2), a, Pt(9), WHITE, b=True)
    txt(s2, x+Inches(0.08), Inches(2.06), Inches(1.5), Inches(0.16), b, Pt(7.5), SUBTLE, f='Consolas')

v_arrow(s2, Inches(6.2), Inches(2.36))

# ── Terraform ──
layer_title(s2, Inches(0.6), Inches(2.58), 'Terraform Layer  (HCL模板 · 7核心变量 · 唯一命名uuid)', GREEN)
box(s2, Inches(0.6), Inches(2.8), Inches(7.5), Inches(0.95), fill=DARK_BG2, border=GREEN)
tag(s2, Inches(0.68), Inches(2.85), 'IaC Template', GREEN)
for i,(a,b,c) in enumerate([('huaweicloud_vpc','172.16.0.0/16',BLUE),('vpc_subnet','172.16.1.0/24',BLUE),
    ('secgroup','SSH:22+Proxy:8787',GREEN),('vpc_eip','300Mbit/s 流量',ORANGE),('compute_instance','Flexus X Ubuntu',PURPLE)]):
    x = Inches(0.68 + (i%3)*2.5)
    y = Inches(3.15 + (i//3)*0.25)
    box(s2, x, y, Inches(2.35), Inches(0.22), r=True, fill=BOX_DARK, border=c)
    txt(s2, x+Inches(0.05), y, Inches(1.05), Inches(0.22), a, Pt(7), c, f='Consolas')
    txt(s2, x+Inches(1.1), y, Inches(1.2), Inches(0.22), b, Pt(7), SUBTLE, f='Consolas')

# user_data
box(s2, Inches(8.3), Inches(2.8), Inches(4.4), Inches(0.95), fill=DARK_BG2, border=GOLD)
tag(s2, Inches(8.38), Inches(2.85), 'user_data Flow', GOLD)
for i, step in enumerate(['① chpasswd 重置密码','② wget OBS 下载 install.sh','③ bash 执行安装脚本','④ rm -rf 清理']):
    txt(s2, Inches(8.38), Inches(3.15 + i*0.14), Inches(4.2), Inches(0.14), step, Pt(8), WHITE)
txt(s2, Inches(8.38), Inches(3.72), Inches(4.2), Inches(0.14), '模板-脚本分离 · user_data最小化', Pt(7), GOLD, f='Consolas')

v_arrow(s2, Inches(6.2), Inches(3.8))

# ── Install Script ──
layer_title(s2, Inches(0.6), Inches(4.05), 'Install Script Layer  (5阶段 · install_*.sh)', ORANGE)
box(s2, Inches(0.6), Inches(4.28), Inches(12.1), Inches(0.55), fill=DARK_BG2, border=ORANGE)
for i,(stg,name,desc,clr) in enumerate([('Stage 1','系统准备','dpkg + apt',GOLD),('Stage 2','Docker','CN华为云/HK官方',BLUE),
    ('Stage 3','应用配置','目录/配置/cron',GREEN),('Stage 4','启动验证','compose up 120s',PURPLE),('Stage 5','OpenCode','Node.js 18+ npm',CYAN)]):
    x = Inches(0.65 + i*2.44)
    box(s2, x, Inches(4.33), Inches(2.35), Inches(0.45), r=True, fill=BOX_DARK, border=clr)
    txt(s2, x+Inches(0.05), Inches(4.35), Inches(2.25), Inches(0.15), f'{stg} · {name}', Pt(8), clr, b=True)
    txt(s2, x+Inches(0.05), Inches(4.52), Inches(2.25), Inches(0.15), desc, Pt(7), SUBTLE, f='Consolas')

v_arrow(s2, Inches(6.2), Inches(4.88))

# ── ECS Runtime ──
layer_title(s2, Inches(0.6), Inches(5.08), 'ECS Runtime  (Flexus X · Ubuntu 24.04 · x1.2u.4g · 40GB SSD)', BLUE)
box(s2, Inches(0.6), Inches(5.3), Inches(12.1), Inches(0.5), fill=DARK_BG2, border=BLUE)
for i,(a,b,c) in enumerate([('OpenCode CLI','npm -g opencode-ai',BLUE),('Headroom Proxy','pip headroom-ai :8787',GREEN),
    ('Context Engine','AST压缩 60-95%',RAG),('MaaS Endpoint','modelarts-maas.com',GOLD),
    ('Docker Services','LiteLLM+PG+Prom',ORANGE),('Monitoring','/metrics /stats',CYAN)]):
    x = Inches(0.65 + i*2.02)
    box(s2, x, Inches(5.34), Inches(1.95), Inches(0.42), r=True, fill=BOX_DARK, border=c)
    txt(s2, x+Inches(0.05), Inches(5.36), Inches(1.85), Inches(0.15), a, Pt(8), c, b=True)
    txt(s2, x+Inches(0.05), Inches(5.53), Inches(1.85), Inches(0.14), b, Pt(7), SUBTLE, f='Consolas')

# ── AI Skills (right) ──
box(s2, Inches(9.5), Inches(1.75), Inches(3.2), Inches(2.0), fill=DARK_BG2, border=GOLD)
txt(s2, Inches(9.6), Inches(1.8), Inches(3), Inches(0.2), 'AI Skills 技能体系', Pt(10), GOLD, b=True)
for i,(a,b) in enumerate([('rfs-practices','创建/修复RFS模板'),('solution-extractor','方案页→结构化Excel'),
    ('page-enhance','方案页面文案增强'),('deep-search-synth','3层深度调研')]):
    y = Inches(2.05 + i*0.4)
    txt(s2, Inches(9.6), y, Inches(3), Inches(0.18), a, Pt(9), WHITE, b=True, f='Consolas')
    txt(s2, Inches(9.6), y+Inches(0.18), Inches(3), Inches(0.14), b, Pt(7.5), SUBTLE)

# ── Scripts (right) ──
box(s2, Inches(9.5), Inches(3.85), Inches(3.2), Inches(1.85), fill=DARK_BG2, border=GREEN)
txt(s2, Inches(9.6), Inches(3.9), Inches(3), Inches(0.2), 'Python/Bash 工具链', Pt(10), GREEN, b=True)
for i,(a,b) in enumerate([('validate_template.py','RFS模板结构验证'),('gen_docx.py','README→Word中英文'),
    ('gen_sac_docx.py','SAC方案报告生成'),('generate_extension.py','.extension配置生成'),
    ('package_solution.sh','方案打包归档')]):
    y = Inches(4.15 + i*0.3)
    txt(s2, Inches(9.6), y, Inches(3), Inches(0.15), a, Pt(8), WHITE, b=True, f='Consolas')
    txt(s2, Inches(9.6), y+Inches(0.15), Inches(3), Inches(0.12), b, Pt(7), SUBTLE)

# ── Region Diff (right) ──
box(s2, Inches(9.5), Inches(5.75), Inches(3.2), Inches(0.8), fill=DARK_BG2, border=CYAN)
txt(s2, Inches(9.6), Inches(5.8), Inches(3), Inches(0.18), 'CN vs HK 差异', Pt(9), CYAN, b=True)
for i,(a,b) in enumerate([('Docker源:','华为云镜像 / 官方'),('SWR:','需要 / 不需要'),('PyPI:','清华源 / 默认')]):
    y = Inches(6.0 + i*0.15)
    txt(s2, Inches(9.6), y, Inches(0.7), Inches(0.14), a, Pt(7), CYAN, b=True, f='Consolas')
    txt(s2, Inches(10.3), y, Inches(2.3), Inches(0.14), b, Pt(7), SUBTLE, f='Consolas')

foot(s2, 'SAC · 详细技术架构', 'Solution Practices')

# ── Save ──
out = r'C:\Users\Administrator\Desktop\Project\claudeproject\solution-practices\ppt\competition\arch_detail.pptx'
prs.save(out)
import os
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}, Size: {os.path.getsize(out)/1024:.0f} KB")
