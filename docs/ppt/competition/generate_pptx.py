#!/usr/bin/env python3
"""Generate competition.pptx - 智云赢商·AI Agent 创新实战 (28 pages)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# Colors
DARK_BG = RGBColor(0x0F, 0x0F, 0x23)
DARK_BG2 = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
SUBTLE_TEXT = RGBColor(0xAA, 0xAA, 0xBB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=Pt(14), color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = sz; p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return box

def mtb(slide, l, t, w, h, lines):
    """Multi-line text box. lines = [(text, size, color, bold, align), ...]"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt, sz, clr, bld, al = item[0], item[1] if len(item)>1 else Pt(14), item[2] if len(item)>2 else DARK_TEXT, item[3] if len(item)>3 else False, item[4] if len(item)>4 else PP_ALIGN.LEFT
        p.text = txt; p.font.size = sz; p.font.color.rgb = clr; p.font.bold = bld; p.font.name = 'Microsoft YaHei'; p.alignment = al
    return box

def chrome(slide, left_text, right_text, dark=True):
    c = WHITE if dark else GRAY_TEXT
    tb(slide, Inches(0.6), Inches(0.25), Inches(6), Inches(0.3), left_text, Pt(10), c, font='Consolas')
    tb(slide, Inches(10), Inches(0.25), Inches(2.8), Inches(0.3), right_text, Pt(10), c, align=PP_ALIGN.RIGHT, font='Consolas')
    rect(slide, Inches(0.6), Inches(0.6), Inches(12.1), Pt(1), fill=RGBColor(0xE2,0xE8,0xF0) if not dark else RGBColor(0x33,0x33,0x50))

def footer(slide, left, right):
    tb(slide, Inches(0.6), Inches(7.0), Inches(6), Inches(0.35), left, Pt(9), GRAY_TEXT, font='Consolas')
    tb(slide, Inches(10), Inches(7.0), Inches(2.8), Inches(0.35), right, Pt(9), GRAY_TEXT, align=PP_ALIGN.RIGHT, font='Consolas')

def kicker(slide, text, left=Inches(0.8), top=Inches(1.0)):
    tb(slide, left, top, Inches(5), Inches(0.35), text, Pt(11), GOLD, bold=True, font='Consolas')

def card(slide, l, t, w, h, title, desc, tag='', color=BLUE):
    rrect(slide, l, t, w, h, fill=WHITE, border=CARD_BORDER)
    tb(slide, l+Inches(0.15), t+Inches(0.1), w-Inches(0.3), Inches(0.35), title, Pt(16), DARK_TEXT, bold=True)
    tb(slide, l+Inches(0.15), t+Inches(0.5), w-Inches(0.3), h-Inches(0.85), desc, Pt(11), GRAY_TEXT)
    if tag:
        tb(slide, l+Inches(0.15), t+h-Inches(0.3), w-Inches(0.3), Inches(0.25), tag, Pt(9), color, font='Consolas')
    return rrect(slide, l, t, w, h, border=CARD_BORDER)

def stat_card(slide, l, t, w, num, unit, desc):
    rrect(slide, l, t, w, Inches(1.6), fill=WHITE, border=CARD_BORDER)
    tb(slide, l+Inches(0.15), t+Inches(0.1), w-Inches(0.3), Inches(0.5), num, Pt(26), BLUE, bold=True)
    tb(slide, l+Inches(0.15), t+Inches(0.6), w-Inches(0.3), Inches(0.3), unit, Pt(12), DARK_TEXT, bold=True)
    tb(slide, l+Inches(0.15), t+Inches(0.95), w-Inches(0.3), Inches(0.5), desc, Pt(9), GRAY_TEXT)

def dark_card(slide, l, t, w, h, title, desc, tag1='', tag2=''):
    rrect(slide, l, t, w, h, fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))
    tb(slide, l+Inches(0.2), t+Inches(0.15), w-Inches(0.4), Inches(0.35), title, Pt(18), WHITE, bold=True)
    tb(slide, l+Inches(0.2), t+Inches(0.55), w-Inches(0.4), h-Inches(1.1), desc, Pt(10), SUBTLE_TEXT)

print("Generating 28 slides...")

# =============== PAGE 1: COVER ===============
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK_BG)
rect(s, Inches(0), Inches(0), Inches(0.12), Inches(7.5), fill=GOLD)
tb(s, Inches(0.8), Inches(1.5), Inches(8), Inches(0.5), '智云赢商 · AI Agent 创新实战大比武', Pt(13), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(2.3), Inches(11), Inches(1.2), 'AI Agent 驱动', Pt(56), WHITE, bold=True)
tb(s, Inches(0.8), Inches(3.4), Inches(11), Inches(0.7), '商业洞察与云方案实践', Pt(38), SUBTLE_TEXT)
rect(s, Inches(0.8), Inches(4.5), Inches(2), Pt(2), fill=GOLD)
tb(s, Inches(0.8), Inches(4.8), Inches(10), Inches(1.2), '基于华为云 MaaS 大模型服务、RFS 资源编排与 Headroom 上下文压缩引擎\n构建智能商业洞察平台与解决方案实践体系', Pt(15), SUBTLE_TEXT)
tb(s, Inches(0.8), Inches(6.4), Inches(5), Inches(0.4), '唐潘 · 解决方案  |  2026', Pt(12), GRAY_TEXT, font='Consolas')
print("  Slide 1: Cover")

# =============== PAGE 2: AGENDA ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, '议程 · Agenda', '02 / 28', dark=False)
kicker(s, 'AGENDA')
tb(s, Inches(0.8), Inches(1.4), Inches(6), Inches(0.6), '分享议程', Pt(38), DARK_TEXT, bold=True)

agenda = [('01', '背景与挑战', '信息孤岛与效率瓶颈'), ('02', '项目全景 · 双轨并行', 'InsightPro + SAC 两大项目'),
          ('03', 'InsightPro · 商业洞察', '让数据说话，让决策有据'), ('04', 'SAC · 解决方案实践', '方案实践，一键部署'),
          ('05', '华为云产品深度使用', 'MaaS · RFS · Headroom'), ('06', '创新性与价值总结', 'AI Agent 驱动业务价值')]
for i, (n, t, d) in enumerate(agenda):
    y = Inches(2.2) + Inches(i * 0.82)
    tb(s, Inches(1.0), y, Inches(0.5), Inches(0.4), n, Pt(18), BLUE, bold=True, font='Consolas')
    tb(s, Inches(1.8), y, Inches(3.5), Inches(0.35), t, Pt(15), DARK_TEXT, bold=True)
    tb(s, Inches(1.8), y+Inches(0.35), Inches(4), Inches(0.3), d, Pt(10), GRAY_TEXT)
# Right side quote
rect(s, Inches(7.5), Inches(2.5), Pt(3), Inches(3.5), fill=BLUE)
tb(s, Inches(8.0), Inches(2.8), Inches(4.5), Inches(3.2), '围绕\n业务价值、技术实现、\n创新性、产品反馈\n\n四大维度深度展示\nAI Agent 解决实际\n业务问题的效果', Pt(14), DARK_TEXT)
footer(s, '智云赢商 · AI Agent 创新实战', 'Agenda')
print("  Slide 2: Agenda")

# =============== PAGE 3: ACT I CURTAIN ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, '第一幕 · 背景与挑战', 'Act I · 03 / 28')
tb(s, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4), 'Act I', Pt(12), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(2.6), Inches(10), Inches(1.2), 'AI Agent 时代', Pt(50), WHITE, bold=True)
tb(s, Inches(0.8), Inches(3.9), Inches(9), Inches(1.0), '从信息过载到智能决策，从重复劳动到自动化开发\nAI Agent 正在重塑我们获取情报和构建解决方案的方式', Pt(18), SUBTLE_TEXT)
footer(s, '第一幕引子 · 为什么需要 AI Agent', '— · —')
print("  Slide 3: Act I Curtain")

# =============== PAGE 4: DUAL CHALLENGES ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, '双重挑战 · Dual Challenges', 'Act I · 04 / 28', dark=False)
kicker(s, 'THE CHALLENGES')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '信息孤岛与效率瓶颈', Pt(36), DARK_TEXT, bold=True)
# Left card
card(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5), '情报获取之困',
     '• 信息分散在邮件/聊天/文档\n• 手动整理Excel，耗时耗力\n• 分析缺乏标准化方法论\n• 知识无法沉淀，人员变动断层', 'Challenge · 商业洞察')
# Right card
card(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.5), '开发部署之痛',
     '• 方案从设计到部署周期漫长\n• 手动配置基础设施，重复劳动\n• AI开发环境搭建复杂，门槛高\n• Token成本消耗大，缺乏优化', 'Challenge · 方案构建')
footer(s, 'Page 04 · 双重挑战', 'Challenges')
print("  Slide 4: Challenges")

# =============== PAGE 5: PROJECT PORTFOLIO ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, '项目全景 · Portfolio', 'Act I · 05 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'PORTFOLIO OVERVIEW', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '双轨并行 · AI Agent 驱动', Pt(36), WHITE, bold=True)
# Card A
rrect(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.2), fill=DARK_BG2, border=RGBColor(0x3B,0x82,0xF6))
tb(s, Inches(1.1), Inches(2.7), Inches(2), Inches(0.3), 'Project A', Pt(10), BLUE, font='Consolas')
tb(s, Inches(1.1), Inches(3.0), Inches(5), Inches(0.5), 'InsightPro', Pt(28), WHITE, bold=True)
tb(s, Inches(1.1), Inches(3.5), Inches(5), Inches(1.0), '企业级AI商业洞察平台\n基于DeepSeek+MaaS双AI引擎\n自动采集多源信息\n竞品追踪/市场分析/招标商机', Pt(12), SUBTLE_TEXT)
tb(s, Inches(1.1), Inches(4.8), Inches(5), Inches(0.25), 'Next.js + FastAPI · DeepSeek + MaaS', Pt(10), GRAY_TEXT, font='Consolas')
# Card B
rrect(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.2), fill=DARK_BG2, border=RGBColor(0x60,0x60,0x80))
tb(s, Inches(7.3), Inches(2.7), Inches(2), Inches(0.3), 'Project B', Pt(10), GRAY_TEXT, font='Consolas')
tb(s, Inches(7.3), Inches(3.0), Inches(5), Inches(0.5), 'SAC · 解决方案实践', Pt(28), WHITE, bold=True)
tb(s, Inches(7.3), Inches(3.5), Inches(5), Inches(1.0), '基于华为云RFS/OpenTofu\n解决方案实践仓库\nHeadroom+OpenCode AI编程\n基础设施即代码扩展', Pt(12), SUBTLE_TEXT)
tb(s, Inches(7.3), Inches(4.8), Inches(5), Inches(0.25), 'RFS + OpenTofu · Headroom + MaaS', Pt(10), GRAY_TEXT, font='Consolas')
# Stats bar
rect(s, Inches(0.8), Inches(6.0), Inches(11.7), Pt(1), fill=RGBColor(0x50,0x50,0x70))
for i, (n, l) in enumerate([('2', 'AI Agent 项目'), ('3', '华为云服务'), ('100%', '开源项目')]):
    tb(s, Inches(1.5+i*4), Inches(6.2), Inches(1.5), Inches(0.4), n, Pt(24), GOLD, bold=True)
    tb(s, Inches(1.5+i*4), Inches(6.6), Inches(2.5), Inches(0.3), l, Pt(11), GRAY_TEXT)
footer(s, 'Page 05 · 双项目定位', 'Portfolio')
print("  Slide 5: Portfolio")

# =============== PAGE 6: ACT II CURTAIN ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, '第二幕 · InsightPro', 'Act II · 06 / 28', dark=False)
tb(s, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4), 'Act II · InsightPro', Pt(12), BLUE, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(2.6), Inches(10), Inches(1.0), '让数据说话', Pt(50), DARK_TEXT, bold=True)
tb(s, Inches(0.8), Inches(3.8), Inches(9), Inches(1.0), '基于 DeepSearcher RAG + 双 AI 引擎的企业级竞争情报分析平台\n将海量信息转化为可执行的商业洞察', Pt(17), GRAY_TEXT)
footer(s, '第二幕 · InsightPro', '— · —')
print("  Slide 6: Act II")

# =============== PAGE 7: BUSINESS VALUE ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, 'InsightPro · 业务价值', 'Act II / Value · 07 / 28', dark=False)
kicker(s, 'BUSINESS VALUE')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '从情报获取到决策支撑', Pt(36), DARK_TEXT, bold=True)
vals = [('6', '大行业', '生物医疗/交通/基建/互联网/零售/制造'),
        ('5', '大厂商', 'AWS/Azure/阿里云/腾讯云/火山云'),
        ('12', '个模块', '从行业全景到招标信息完整体系'),
        ('24h', '实时监控', '百度热搜+GitHub Trending'),
        ('16', '个端点', 'FastAPI RESTful API'),
        ('94.74.90.21:4000', '已部署上线', '持续迭代中')]
for i, (n, u, d) in enumerate(vals):
    row, col = i // 3, i % 3
    x, y = Inches(0.8 + col * 4.1), Inches(2.3 + row * 2.2)
    stat_card(s, x, y, Inches(3.6), n, u, d)
footer(s, 'InsightPro · 业务价值量化', 'Act II · Value')
print("  Slide 7: Value")

# =============== PAGE 8: ARCHITECTURE ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, 'InsightPro · 技术架构', 'Act II / Architecture · 08 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'SYSTEM ARCHITECTURE', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '三层智能架构', Pt(36), WHITE, bold=True)

layers = [
    ('Layer 01 · 数据采集层', [
        ('百度热搜', 'BeautifulSoup抓取'), ('GitHub Trend', '每日09:00刷新'),
        ('竞品数据', '5大云厂商'), ('SQLite缓存', '失败降级机制')]),
    ('Layer 02 · AI分析层（华为云MaaS）', [
        ('DeepSeek API', '结构化分析引擎'), ('MaaS 服务', 'ModelArts流式对话'),
        ('APScheduler', '每日自动分析'), ('降级机制', '不可用时返回缓存')]),
    ('Layer 03 · 展示交互层', [
        ('Next.js 14', 'App Router 16路由'), ('FastAPI', '9个RESTful端点'),
        ('Recharts', '雷达/柱状/趋势图'), ('Tailwind CSS', '响应式布局')]),
]
for li, (layer_name, items) in enumerate(layers):
    y = Inches(2.3 + li * 1.7)
    tb(s, Inches(0.8), y, Inches(8), Inches(0.3), layer_name, Pt(11), GOLD, font='Consolas')
    for ii, (t, d) in enumerate(items):
        x = Inches(0.8 + ii * 3.1)
        rrect(s, x, y+Inches(0.4), Inches(2.8), Inches(1.0), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))
        tb(s, x+Inches(0.15), y+Inches(0.5), Inches(2.5), Inches(0.3), t, Pt(13), WHITE, bold=True)
        tb(s, x+Inches(0.15), y+Inches(0.85), Inches(2.5), Inches(0.3), d, Pt(9), GRAY_TEXT)
footer(s, 'InsightPro · 技术实现', 'Architecture')
print("  Slide 8: Architecture")

# =============== PAGE 9: 12 FEATURES ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, 'InsightPro · 核心功能', 'Act II / Features · 09 / 28', dark=False)
kicker(s, 'CORE CAPABILITIES')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '十二大功能模块', Pt(36), DARK_TEXT, bold=True)

modules = [('今日洞察','首页总览'),('行业全景','六大行业动态'),('案例库','上云最佳实践'),
           ('热点追踪','GitHub Trend'),('友商洞察','6大场景化分析'),('政策法规','核心政策追踪'),
           ('商业快讯','全球媒体头条'),('增长机会','客户切入策略'),('市场情报','商机信息分析'),
           ('需求挖掘','客户需求线索'),('招标信息','招标情报'),('数据大屏','运营指标图表')]
for i, (n, d) in enumerate(modules):
    row, col = i // 4, i % 4
    x, y = Inches(0.8 + col * 3.1), Inches(2.3 + row * 1.5)
    rrect(s, x, y, Inches(2.8), Inches(1.25), fill=WHITE, border=CARD_BORDER)
    tb(s, x+Inches(0.2), y+Inches(0.2), Inches(2.4), Inches(0.35), n, Pt(14), DARK_TEXT, bold=True)
    tb(s, x+Inches(0.2), y+Inches(0.65), Inches(2.4), Inches(0.35), d, Pt(10), GRAY_TEXT)
footer(s, 'InsightPro · 功能全景', 'Features')
print("  Slide 9: Features")

# =============== PAGE 10: COMPETITIVE SCENES ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, 'InsightPro · 竞争分析', 'Act II / Deep Dive · 10 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'COMPETITIVE LANDSCAPE', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '六维竞争场景', Pt(36), WHITE, bold=True)
tb(s, Inches(0.8), Inches(2.1), Inches(7), Inches(0.6), '华为云 vs AWS / Azure / 阿里云 / 腾讯云 / 火山云', Pt(14), SUBTLE_TEXT)

scenes = [('场景 01','中资企业出海','vs AWS + Azure'),('场景 02','AI大模型开放平台','vs Azure + AWS'),
          ('场景 03','开发者生态争夺','vs AWS+Azure+阿里云'),('场景 04','企业SaaS生态','vs Azure + 阿里云'),
          ('场景 05','腰部客户性价比','vs 阿里云 + 火山云'),('场景 06','传统行业数字化','vs 阿里云 + 腾讯云')]
for i, (n, t, v) in enumerate(scenes):
    row, col = i // 3, i % 3
    x, y = Inches(0.8 + col * 4.1), Inches(2.8 + row * 1.8)
    rrect(s, x, y, Inches(3.8), Inches(1.5), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))
    tb(s, x+Inches(0.2), y+Inches(0.15), Inches(2), Inches(0.25), n, Pt(10), GOLD, font='Consolas')
    tb(s, x+Inches(0.2), y+Inches(0.45), Inches(3.4), Inches(0.4), t, Pt(15), WHITE, bold=True)
    tb(s, x+Inches(0.2), y+Inches(1.0), Inches(3.4), Inches(0.3), v, Pt(10), GRAY_TEXT)
footer(s, 'InsightPro · 竞争分析深度', 'Deep Dive')
print("  Slide 10: Competitive")

# =============== PAGE 11: DUAL ENGINE ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, 'InsightPro · 华为云 MaaS 集成', 'Act II / MaaS · 11 / 28', dark=False)
kicker(s, 'HUAWEI CLOUD INTEGRATION')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), 'DeepSeek + MaaS 双引擎', Pt(36), DARK_TEXT, bold=True)
card(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.0), 'Engine A · DeepSeek API',
     '结构化分析引擎\n输入市场数据或竞品信息\n输出多维度深度分析报告\n通过OpenAI兼容SDK集成', '报告生成 · 洞察合成', BLUE)
card(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.0), 'Engine B · ModelArts MaaS',
     '流式对话引擎\n基于华为云MaaS部署deepseek-v3.2\n支持实时问答与数据查询\n流式响应面向终端用户', '实时问答 · 流式响应', BLUE)
for i, (n, d) in enumerate([('OpenAI SDK','API兼容·零迁移'),('DeepSeek-V3','高性价比推理'),('双模式','结构分析+流式对话')]):
    tb(s, Inches(2+i*4), Inches(5.8), Inches(2.5), Inches(0.35), n, Pt(18), BLUE, bold=True)
    tb(s, Inches(2+i*4), Inches(6.15), Inches(3), Inches(0.3), d, Pt(10), GRAY_TEXT)
footer(s, 'InsightPro · 华为云 MaaS', 'MaaS Integration')
print("  Slide 11: Dual Engine")

# =============== PAGE 12: FEEDBACK ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, 'InsightPro · 产品反馈', 'Act II / Feedback · 12 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'PRODUCT FEEDBACK', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '从使用中学习，在反馈中迭代', Pt(34), WHITE, bold=True)

fb = [('信息聚合效率', '百度热搜+GitHub Trending自动采集\n替代人工搜集，效率提升80%+\n多源数据统一展示', '效率提升 80%'),
      ('竞争分析深度', '6大竞争场景结构化分析框架\n差距量化+策略建议+时间窗口\n从Excel到AI驱动标准体系', '结构化输出'),
      ('持续迭代方向', 'Supabase持久化/AI结果入库\n用户认证/全局搜索/通知系统\nPDF导出/移动端适配进行中', '持续演进')]
for i, (t, d, tag) in enumerate(fb):
    x = Inches(0.8 + i * 4.1)
    dark_card(s, x, Inches(2.5), Inches(3.8), Inches(3.2), t, d)
    tb(s, x+Inches(0.25), Inches(5.3), Inches(3.3), Inches(0.25), tag, Pt(14), GOLD, bold=True)
# Tags
for i, t in enumerate(['MVP上线','12模块','开源','降级机制']):
    tb(s, Inches(1.5+i*3), Inches(6.3), Inches(2.5), Inches(0.3), t, Pt(11), GRAY_TEXT, font='Consolas')
footer(s, 'InsightPro · 反馈与迭代', 'Feedback')
print("  Slide 12: Feedback")

# =============== PAGE 13: ACT III CURTAIN ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, '第三幕 · SAC 解决方案实践', 'Act III · 13 / 28', dark=False)
tb(s, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4), 'Act III · Solution Practices', Pt(12), BLUE, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(2.6), Inches(10), Inches(1.0), '方案实践', Pt(50), DARK_TEXT, bold=True)
tb(s, Inches(0.8), Inches(3.8), Inches(9), Inches(1.0), '基于华为云 RFS/OpenTofu\n将解决方案从设计、部署到运维以代码形式标准化、自动化、可复用', Pt(17), GRAY_TEXT)
footer(s, '第三幕 · SAC', '— · —')
print("  Slide 13: Act III")

# =============== PAGE 14: SAC PORTFOLIO ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, 'SAC · 解决方案全景', 'Act III / Portfolio · 14 / 28', dark=False)
kicker(s, 'SOLUTION PORTFOLIO')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '六大解决方案实践', Pt(36), DARK_TEXT, bold=True)

sols = [('Headroom\nOpenCode','AI编程助手平台\nToken压缩60-95%','CN/HK'),
        ('LiteLLM','多模型API网关\n统一接入多种大模型','CN/HK'),
        ('OpenHands','AI驱动的\n开发助手平台','CN'),
        ('Supabase','开源Firebase替代\n后端即服务','CN'),
        ('CodeWhale','智能代码\n分析平台','CN'),
        ('AiToEarn','AI收益系统','CN/HK')]
for i, (t, d, r) in enumerate(sols):
    row, col = i // 3, i % 3
    x, y = Inches(0.8 + col * 4.1), Inches(2.3 + row * 2.1)
    rrect(s, x, y, Inches(3.8), Inches(1.85), fill=WHITE, border=CARD_BORDER)
    tb(s, x+Inches(0.2), y+Inches(0.15), Inches(3.4), Inches(0.65), t, Pt(16), DARK_TEXT, bold=True)
    tb(s, x+Inches(0.2), y+Inches(0.85), Inches(3.4), Inches(0.6), d, Pt(11), GRAY_TEXT)
    tb(s, x+Inches(0.2), y+Inches(1.45), Inches(3.4), Inches(0.25), r, Pt(10), BLUE, font='Consolas')

rect(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.3), fill=RGBColor(0xF0,0xF4,0xF8))
tb(s, Inches(1.0), Inches(6.72), Inches(11), Inches(0.25), '统一目录结构：docs · scripts · terraform  |  标准化 · 自动化 · 可复用', Pt(10), GRAY_TEXT, align=PP_ALIGN.CENTER)
footer(s, 'SAC · 解决方案全景', 'Portfolio')
print("  Slide 14: SAC Portfolio")

# =============== PAGE 15: HEADROOM CORE ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, 'SAC · Headroom OpenCode', 'Act III / Headroom · 15 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'THE FRUGAL CODING ASSISTANT', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), 'Headroom + OpenCode', Pt(36), WHITE, bold=True)
tb(s, Inches(0.8), Inches(2.1), Inches(8), Inches(0.6), '开源 AI 编程环境，通过智能上下文压缩引擎将 Token 成本降低 60-95%', Pt(14), SUBTLE_TEXT)

vals4 = [('60-95%','Token压缩','智能分析工具/文件/历史',GOLD),('零费用','开源·MIT','无供应商锁定',BLUE),
         ('一键部署','Flexus X','RFS模板3-5分钟',GREEN),('多模型','DeepSeek/Claude','灵活切换适配场景',RGBColor(0xA7,0x8B,0xF0))]
for i, (n, t, d, c) in enumerate(vals4):
    x = Inches(0.8 + i * 3.1)
    rrect(s, x, Inches(3.0), Inches(2.8), Inches(2.5), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))
    tb(s, x+Inches(0.2), Inches(3.2), Inches(2.4), Inches(0.5), n, Pt(24), c, bold=True)
    tb(s, x+Inches(0.2), Inches(3.7), Inches(2.4), Inches(0.3), t, Pt(14), WHITE, bold=True)
    tb(s, x+Inches(0.2), Inches(4.2), Inches(2.4), Inches(0.5), d, Pt(10), GRAY_TEXT)
footer(s, 'Headroom OpenCode · 核心方案', 'Core Solution')
print("  Slide 15: Headroom")

# =============== PAGE 16: HEADROOM ARCH ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, 'SAC · Headroom 技术架构', 'Act III / Architecture · 16 / 28', dark=False)
kicker(s, 'ARCHITECTURE')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '系统架构与数据流', Pt(36), DARK_TEXT, bold=True)

flow = [('开发者\nOpenCode CLI',BLUE),('Headroom Proxy\nport 8787',RGBColor(0x3B,0x82,0xF6)),
        ('华为云MaaS API',GREEN),('LLM推理\nDeepSeek/Claude',RGBColor(0xA7,0x8B,0xF0))]
for i, (lab, c) in enumerate(flow):
    x = Inches(0.8 + i * 3.1)
    rrect(s, x, Inches(2.5), Inches(2.8), Inches(1.5), fill=WHITE, border=c)
    tb(s, x+Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.8), lab, Pt(14), DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        tb(s, x+Inches(2.6), Inches(2.9), Inches(0.6), Inches(0.4), '→', Pt(20), GRAY_TEXT, align=PP_ALIGN.CENTER)

# Tech note
rrect(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.0), fill=RGBColor(0xEE,0xF2,0xFF), border=BLUE)
tb(s, Inches(1.1), Inches(4.6), Inches(5), Inches(0.25), '上下文压缩引擎', Pt(10), BLUE, font='Consolas')
tb(s, Inches(1.1), Inches(4.9), Inches(11), Inches(0.4), 'AST感知压缩 · 可逆压缩(按需解压) · Prometheus监控 · 数据本地处理 · 多模型支持', Pt(12), DARK_TEXT)

for i, (n, d) in enumerate([('Flexus X ECS','x1.2u.4g'),('VPC+Subnet','172.16.0.0/16'),('安全组','22+8787'),('弹性IP','300Mbit/s')]):
    tb(s, Inches(0.8+i*3.1), Inches(5.8), Inches(2.5), Inches(0.3), n, Pt(13), DARK_TEXT, bold=True)
    tb(s, Inches(0.8+i*3.1), Inches(6.1), Inches(2.5), Inches(0.3), d, Pt(10), GRAY_TEXT)
footer(s, 'SAC · 技术架构', 'Architecture')
print("  Slide 16: Headroom Arch")

# =============== PAGE 17: RFS ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, 'SAC · 华为云 RFS', 'Act III / RFS · 17 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'RFS RESOURCE ORCHESTRATION', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '基础设施即代码 · 一键部署', Pt(34), WHITE, bold=True)

# Left
rrect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(3.5), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))
tb(s, Inches(1.1), Inches(2.5), Inches(4), Inches(0.25), '模板结构', Pt(10), GOLD, font='Consolas')
tb(s, Inches(1.1), Inches(2.9), Inches(5), Inches(2.0), 'Variables → Locals → Data Sources\n→ VPC/Subnet/SecGroup\n→ EIP → ECS\n→ user_data 下载脚本\n→ 执行 → 清理', Pt(14), WHITE, font='Consolas')
tb(s, Inches(1.1), Inches(5.0), Inches(3), Inches(0.25), '模板-脚本分离原则', Pt(10), GOLD, font='Consolas')

# Right
rrect(s, Inches(7.0), Inches(2.3), Inches(5.5), Inches(3.5), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))
tb(s, Inches(7.3), Inches(2.5), Inches(4), Inches(0.25), 'RFS 一键部署流程', Pt(10), GOLD, font='Consolas')
for i, step in enumerate(['1. 打开 RFS 部署链接','2. 填写参数配置','3. 点击"一键部署"','4. 自动创建运行','5. SSH 验证服务']):
    tb(s, Inches(7.3), Inches(2.9+i*0.5), Inches(4.5), Inches(0.35), step, Pt(13), WHITE)

# Bottom info
rrect(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8), fill=DARK_BG2, border=RGBColor(0x3B,0x82,0xF6))
tb(s, Inches(1.1), Inches(6.1), Inches(11), Inches(0.5), '.tf user_data: 重置密码→OBS下载→执行→清理  |  .sh安装脚本: Docker安装→容器启动→健康检查', Pt(12), WHITE)
footer(s, 'SAC · RFS 资源编排', 'RFS')
print("  Slide 17: RFS")

# =============== PAGE 18: MATRIX ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, 'SAC · 多方案矩阵', 'Act III / Matrix · 18 / 28', dark=False)
kicker(s, 'SOLUTION MATRIX')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '方案场景与价值矩阵', Pt(34), DARK_TEXT, bold=True)

headers = ['方案名称', '描述', '技术栈', '区域']
widths = [Inches(2.5), Inches(4.5), Inches(3.5), Inches(1.2)]
x0 = Inches(0.8)
y0 = Inches(2.2)
# Header
for i, (h, w) in enumerate(zip(headers, widths)):
    x = x0 + sum(widths[:i])
    rect(s, x, y0, w, Inches(0.4), fill=DARK_BG)
    tb(s, x+Inches(0.1), y0+Inches(0.02), w-Inches(0.2), Inches(0.35), h, Pt(11), WHITE, bold=True)

rows = [('Headroom OpenCode','AI编程助手·Token压缩60-95%','Flexus X·一键部署','CN/HK'),
        ('LiteLLM','多模型API网关·统一接入','MaaS兼容·负载均衡','CN/HK'),
        ('OpenHands','AI开发助手·代码生成','MaaS后端·沙箱执行','CN'),
        ('Supabase','开源Firebase·后端即服务','PostgreSQL·实时订阅','CN'),
        ('CodeWhale','智能代码分析平台','代码审查·质量分析','CN'),
        ('AiToEarn','AI收益系统','自动化·智能调度','CN/HK')]
for ri, (c1, c2, c3, c4) in enumerate(rows):
    y = y0 + Inches(0.4 + ri * 0.5)
    bg = WHITE if ri%2==0 else RGBColor(0xF8,0xFA,0xFC)
    for ci, (val, w) in enumerate(zip([c1,c2,c3,c4], widths)):
        x = x0 + sum(widths[:ci])
        rect(s, x, y, w, Inches(0.5), fill=bg)
        tb(s, x+Inches(0.1), y+Inches(0.05), w-Inches(0.2), Inches(0.4), val, Pt(11), DARK_TEXT, bold=(ci==0))
footer(s, 'SAC · 方案矩阵', 'Matrix')
print("  Slide 18: Matrix")

# =============== PAGE 19: SAC FEEDBACK ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, 'SAC · 产品反馈', 'Act III / Feedback · 19 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'FEEDBACK & SCENARIO DEPTH', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '实际使用场景与反馈', Pt(34), WHITE, bold=True)

sf = [('日常开发辅助','代码补全/重构/Bug修复\n压缩多轮对话重复上下文\n节省60-80%不影响代码质量','节省60-80%'),
      ('大代码库探索','大量文件读取场景\nAST感知压缩率88%\n保留代码结构和语义','压缩85-92%'),
      ('日志调试分析','调试大量日志文件\n最高压缩率92%\nAI分析日志极富性价比','压缩90-95%')]
for i, (t, d, tag) in enumerate(sf):
    x = Inches(0.8 + i * 4.1)
    dark_card(s, x, Inches(2.5), Inches(3.8), Inches(2.8), t, d)
    tb(s, x+Inches(0.25), Inches(5.0), Inches(3.3), Inches(0.25), tag, Pt(16), GOLD, bold=True)

for i, t in enumerate(['月费~$36.50','开源MIT','数据主权','Prometheus']):
    tb(s, Inches(1.5+i*3), Inches(5.8), Inches(2.5), Inches(0.3), t, Pt(11), GRAY_TEXT, font='Consolas')
footer(s, 'SAC · 反馈与场景', 'Feedback')
print("  Slide 19: SAC Feedback")

# =============== PAGE 20: ACT IV CURTAIN ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, '第四幕 · 华为云产品深度', 'Act IV · 20 / 28')
tb(s, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4), 'Act IV · Huawei Cloud Deep Dive', Pt(12), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(2.6), Inches(10), Inches(1.0), '云上深度集成', Pt(50), WHITE, bold=True)
tb(s, Inches(0.8), Inches(3.8), Inches(9), Inches(1.0), 'MaaS 大模型服务、RFS 资源编排、Headroom 上下文压缩\n三项华为云产品在两个项目中深度应用', Pt(17), SUBTLE_TEXT)
footer(s, '第四幕 · 华为云产品深度使用', '— · —')
print("  Slide 20: Act IV")

# =============== PAGE 21: MAAS ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, '华为云 MaaS', 'Act IV / MaaS · 21 / 28', dark=False)
kicker(s, 'MODELARTS MAAS')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '大模型即服务 · 双项目全覆盖', Pt(34), DARK_TEXT, bold=True)

card(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.0), 'InsightPro 中的 MaaS',
     '• 双AI引擎：DeepSeek API + MaaS\n• 流式响应：面向终端实时问答\n• 结构化报告：深度分析报告生成\n• OpenAI兼容SDK：零迁移成本', '流式对话 + 结构化分析', BLUE)
card(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.0), 'SAC 中的 MaaS',
     '• Headroom压缩后转MaaS API\n• 支持DeepSeek/Claude切换\n• MaaS作为统一推理入口\n• 按Token用量计费，成本可控', 'AI 编程推理后端', BLUE)

for i, (n, d) in enumerate([('双项目','InsightPro+SAC全覆盖'),('DeepSeek-V3','高性价比推理'),('OpenAI SDK','标准API兼容')]):
    tb(s, Inches(2+i*4), Inches(5.9), Inches(2.5), Inches(0.35), n, Pt(20), BLUE, bold=True)
    tb(s, Inches(2+i*4), Inches(6.25), Inches(3), Inches(0.3), d, Pt(10), GRAY_TEXT)
footer(s, '华为云 MaaS · 深度使用', 'MaaS')
print("  Slide 21: MaaS")

# =============== PAGE 22: RFS PROCESS ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, '华为云 RFS', 'Act IV / RFS · 22 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'RESOURCE FORMATION SERVICE', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), 'RFS 模板开发与部署', Pt(34), WHITE, bold=True)

steps = [('01','需求分析','资源规划与参数设计'),('02','Terraform','编写.tf模板文件'),
         ('03','安装脚本','独立.sh部署逻辑'),('04','OBS上传','模板+脚本托管'),('05','RFS部署','一键部署·验证服务')]
for i, (n, t, d) in enumerate(steps):
    x = Inches(0.8 + i * 2.45)
    rrect(s, x, Inches(2.5), Inches(2.2), Inches(2.0), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))
    tb(s, x+Inches(0.15), Inches(2.7), Inches(1.9), Inches(0.35), n, Pt(22), GOLD, bold=True, font='Consolas')
    tb(s, x+Inches(0.15), Inches(3.1), Inches(1.9), Inches(0.35), t, Pt(15), WHITE, bold=True)
    tb(s, x+Inches(0.15), Inches(3.6), Inches(1.9), Inches(0.5), d, Pt(10), GRAY_TEXT)
    if i < 4:
        tb(s, x+Inches(2.0), Inches(3.2), Inches(0.5), Inches(0.35), '→', Pt(18), GRAY_TEXT)

rrect(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.2), fill=DARK_BG2, border=RGBColor(0x3B,0x82,0xF6))
tb(s, Inches(1.1), Inches(5.2), Inches(5), Inches(0.25), '模板-脚本分离原则', Pt(10), GOLD, font='Consolas')
tb(s, Inches(1.1), Inches(5.5), Inches(11), Inches(0.5), '.tf user_data仅做：重置密码→OBS下载→执行→清理\n.sh安装脚本包含所有部署逻辑：Docker安装、容器启动、健康检查', Pt(12), WHITE)
footer(s, '华为云 RFS · 资源编排', 'RFS')
print("  Slide 22: RFS Process")

# =============== PAGE 23: HEADROOM MATRIX ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, 'Headroom 压缩引擎', 'Act IV / Headroom · 23 / 28', dark=False)
kicker(s, 'CONTEXT COMPRESSION ENGINE')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '智能压缩 · 大幅降低 Token 成本', Pt(34), DARK_TEXT, bold=True)

comps = [('日常开发','60-80%','代码补全/重构/Bug修复'),('大代码库探索','85-92%','AST感知/保留代码结构'),
         ('架构设计','70-85%','复杂推理/系统设计'),('日志分析调试','90-95%','日志内容/最高压缩率')]
for i, (sc, rate, desc) in enumerate(comps):
    row, col = i // 2, i % 2
    x, y = Inches(0.8 + col * 6.2), Inches(2.3 + row * 1.4)
    rrect(s, x, y, Inches(5.8), Inches(1.2), fill=WHITE, border=CARD_BORDER)
    tb(s, x+Inches(0.2), y+Inches(0.15), Inches(2.5), Inches(0.35), sc, Pt(16), DARK_TEXT, bold=True)
    tb(s, x+Inches(3.5), y+Inches(0.1), Inches(2), Inches(0.4), rate, Pt(22), BLUE, bold=True)
    tb(s, x+Inches(0.2), y+Inches(0.6), Inches(5.4), Inches(0.4), desc, Pt(11), GRAY_TEXT)

rect(s, Inches(0.8), Inches(5.3), Inches(11.7), Pt(1), fill=CARD_BORDER)
tb(s, Inches(0.8), Inches(5.4), Inches(5), Inches(0.3), '技术亮点', Pt(12), DARK_TEXT, bold=True)
hls = ['AST多语言','可逆压缩','Prometheus','数据本地','多模型']
for i, h in enumerate(hls):
    rrect(s, Inches(0.8+i*2.5), Inches(5.8), Inches(2.2), Inches(0.5), fill=RGBColor(0xEE,0xF2,0xFF), border=BLUE)
    tb(s, Inches(0.9+i*2.5), Inches(5.85), Inches(2.0), Inches(0.4), h, Pt(10), BLUE, align=PP_ALIGN.CENTER)
footer(s, 'Headroom · 上下文压缩引擎', 'Headroom')
print("  Slide 23: Headroom Matrix")

# =============== PAGE 24: ACT V CURTAIN ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, '第五幕 · 创新与价值', 'Act V · 24 / 28', dark=False)
tb(s, Inches(0.8), Inches(2.0), Inches(5), Inches(0.4), 'Act V · Innovation & Value', Pt(12), BLUE, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(2.6), Inches(10), Inches(1.0), '创新驱动价值', Pt(50), DARK_TEXT, bold=True)
tb(s, Inches(0.8), Inches(3.8), Inches(9), Inches(1.0), 'AI Agent 不只是技术工具，而是重塑业务流程的驱动力\n从情报获取到方案交付，每一个环节都在被重新定义', Pt(17), GRAY_TEXT)
footer(s, '第五幕 · 创新性与价值总结', '— · —')
print("  Slide 24: Act V")

# =============== PAGE 25: INNOVATION ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, 'AI Agent · 创新实践', 'Act V / Innovation · 25 / 28', dark=False)
kicker(s, 'INNOVATION HIGHLIGHTS')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), 'AI Agent 应用创新点', Pt(34), DARK_TEXT, bold=True)

innos = [('双 AI 引擎协同','InsightPro首创DeepSeek结构化分析+MaaS流式对话双引擎\n不同任务匹配不同引擎，成本与质量平衡优化'),
         ('上下文压缩革命','Headroom AST感知压缩使Token成本降低90%+\n可逆压缩机制按需还原，不牺牲回答质量'),
         ('解决方案实践','将IaC理念扩展到完整解决方案生命周期\nSAC框架使设计/部署/运维全流程代码化标准化'),
         ('Agent辅助竞争分析','AI驱动六维竞争场景分析框架\n每个场景量化差距、提供策略建议'),
         ('自动化情报流水线','从多源采集到AI分析到可视化展示全流程自动化\n定时任务+降级机制确保高可用'),
         ('Skill技能市场','SAC框架内建AI技能体系\nrfs-practices/extractor/page-enhance等实现开发自动化')]
for i, (t, d) in enumerate(innos):
    row, col = i // 3, i % 3
    x, y = Inches(0.8 + col * 4.1), Inches(2.3 + row * 2.2)
    rrect(s, x, y, Inches(3.8), Inches(2.0), fill=WHITE, border=CARD_BORDER)
    tb(s, x+Inches(0.2), y+Inches(0.15), Inches(3.4), Inches(0.35), t, Pt(14), DARK_TEXT, bold=True)
    tb(s, x+Inches(0.2), y+Inches(0.6), Inches(3.4), Inches(1.2), d, Pt(10), GRAY_TEXT)
footer(s, 'AI Agent · 创新实践', 'Innovation')
print("  Slide 25: Innovation")

# =============== PAGE 26: VALUE SUMMARY ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, '业务价值 · 总结', 'Act V / Value · 26 / 28')
tb(s, Inches(0.8), Inches(1.0), Inches(5), Inches(0.35), 'VALUE SUMMARY', Pt(11), GOLD, bold=True, font='Consolas')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '业务价值量化', Pt(34), WHITE, bold=True)

vsum = [('80%','情报效率提升','信息聚合替代人工搜集'),('60-95%','Token节省','上下文压缩降低AI成本'),
        ('一键','部署效率','RFS模板3-5分钟'),('2个','开源项目','全部代码可复用'),
        ('3项','华为云服务','MaaS·RFS·Headroom'),('6+6+4','场景覆盖','竞争+部署+开发全覆盖')]
for i, (n, t, d) in enumerate(vsum):
    row, col = i // 3, i % 3
    x, y = Inches(0.8 + col * 4.1), Inches(2.3 + row * 2.2)
    rrect(s, x, y, Inches(3.8), Inches(2.0), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60))
    tb(s, x+Inches(0.2), y+Inches(0.2), Inches(3.4), Inches(0.5), n, Pt(28), GOLD, bold=True)
    tb(s, x+Inches(0.2), y+Inches(0.8), Inches(3.4), Inches(0.3), t, Pt(15), WHITE, bold=True)
    tb(s, x+Inches(0.2), y+Inches(1.2), Inches(3.4), Inches(0.5), d, Pt(10), GRAY_TEXT)
footer(s, 'Page 26 · 业务价值量化', 'Value')
print("  Slide 26: Value Summary")

# =============== PAGE 27: ROADMAP ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT_BG)
chrome(s, '未来规划 · Roadmap', 'Act V / Future · 27 / 28', dark=False)
kicker(s, 'ROADMAP')
tb(s, Inches(0.8), Inches(1.4), Inches(8), Inches(0.6), '持续演进 · 未来方向', Pt(34), DARK_TEXT, bold=True)

# Left - InsightPro
card(s, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5), 'InsightPro 规划',
     '• Supabase持久化(替换SQLite)\n• AI分析结果入库与历史对比\n• 用户认证(NextAuth+Supabase)\n• 全局搜索功能\n• 通知系统/PDF导出/移动端\n• 更多实时数据源接入',
     '持续进化', BLUE)
# Right - SAC
card(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.5), 'SAC 规划',
     '• 更多解决方案场景覆盖\n• AI Agent自动化程度提升\n• 跨方案编排与组合\n• 方案质量自动评估\n• 社区贡献与协作机制\n• 多区域部署模板扩展',
     '持续进化', BLUE)

for i, t in enumerate(['持续迭代','社区驱动','云原生','AI First']):
    tb(s, Inches(2+i*3), Inches(6.3), Inches(2.5), Inches(0.3), t, Pt(12), BLUE, font='Consolas', align=PP_ALIGN.CENTER)
footer(s, '未来规划 · 持续演进', 'Roadmap')
print("  Slide 27: Roadmap")

# =============== PAGE 28: CLOSING ===============
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, DARK_BG)
chrome(s, 'The Takeaway · 核心洞察', '28 / 28')
tb(s, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4), '智云赢商 · AI Agent 创新实战', Pt(12), GOLD, bold=True, font='Consolas')

# Core quotes
tb(s, Inches(0.8), Inches(2.5), Inches(11), Inches(0.7), 'AI Agent 不是替代人，', Pt(38), WHITE, bold=True)
tb(s, Inches(0.8), Inches(3.3), Inches(11), Inches(0.7), '而是放大人的能力，', Pt(38), WHITE, bold=True)
tb(s, Inches(0.8), Inches(4.1), Inches(11), Inches(0.7), '让每一次决策都有数据支撑。', Pt(38), GOLD, bold=True)
rect(s, Inches(0.8), Inches(4.9), Inches(3), Pt(1), fill=GOLD)

tb(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5), 'InsightPro 让商业洞察智能化，SAC 让方案交付自动化\n基于华为云，用 AI Agent 驱动业务价值落地', Pt(15), SUBTLE_TEXT)
tb(s, Inches(0.8), Inches(5.8), Inches(11), Inches(0.3), 'github.com/Justin-TangPan/Insightpro', Pt(12), BLUE, font='Consolas')
tb(s, Inches(0.8), Inches(6.1), Inches(11), Inches(0.3), 'github.com/Justin-TangPan/solution-practices', Pt(12), BLUE, font='Consolas')
tb(s, Inches(0.8), Inches(6.5), Inches(5), Inches(0.3), '唐潘 · 解决方案  |  2026', Pt(12), GRAY_TEXT, font='Consolas')
footer(s, '谢谢 · 欢迎交流', '— 2026 —')
print("  Slide 28: Closing")

# =============== SAVE ===============
output = r'C:\Users\Administrator\Desktop\Project\claudeproject\solution-practices\ppt\competition\competition.pptx'
prs.save(output)
print(f"\n✅ Saved: {output}")
print(f"   Total: {len(prs.slides)} slides")
import os
print(f"   Size: {os.path.getsize(output)/1024:.1f} KB")
