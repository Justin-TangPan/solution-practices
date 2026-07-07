#!/usr/bin/env python3
"""Generate 智云赢商·AI Agent创新实战 PPT (28 pages, modern magazine style)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x0F, 0x23)   # deep navy dark
DARK_BG2      = RGBColor(0x1A, 0x1A, 0x2E)   # slightly lighter dark
ACCENT_GOLD   = RGBColor(0xF5, 0xA6, 0x23)   # gold accent
ACCENT_BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # bright blue
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)   # cyan
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xD1, 0xD5, 0xDB)
MID_GRAY      = RGBColor(0x9C, 0xA3, 0xAF)
DARK_TEXT      = RGBColor(0x1F, 0x29, 0x37)
LIGHT_BG      = RGBColor(0xF8, 0xFA, 0xFC)   # off-white
CARD_BG       = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER   = RGBColor(0xE5, 0xE7, 0xEB)
GREEN_ACCENT  = RGBColor(0x10, 0xB9, 0x81)
RED_ACCENT    = RGBColor(0xEF, 0x44, 0x44)
PURPLE_ACCENT = RGBColor(0x8B, 0x5C, 0xF6)

# ── Font settings ──────────────────────────────────────────────
TITLE_FONT = "Times New Roman"  # serif
BODY_FONT  = "Microsoft YaHei"

# ── Helper functions ───────────────────────────────────────────

def add_bg(slide, color):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_name=BODY_FONT, font_size=14,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    # set East-Asian font
    from pptx.oxml.ns import qn
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), font_name)
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_name=BODY_FONT, font_size=14,
                   color=WHITE, alignment=PP_ALIGN.LEFT, line_spacing=1.3, anchor=MSO_ANCHOR.TOP):
    """Add a text box with multiple lines, each line different properties possible."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, sz, bld, clr = item, font_size, False, color
        else:
            txt = item.get("text", "")
            sz = item.get("size", font_size)
            bld = item.get("bold", False)
            clr = item.get("color", color)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.name = font_name
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = clr
        p.alignment = alignment
        p.line_spacing = Pt(sz * line_spacing)
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        from pptx.oxml.ns import qn
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:ea'), font_name)
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_GOLD, height=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_data_card(slide, left, top, width, height, number, label, num_color=ACCENT_GOLD, bg_color=CARD_BG):
    """Add a data card with big number and label."""
    # card background
    r = add_shape(slide, left, top, width, height, fill_color=bg_color, line_color=CARD_BORDER)
    # corner radius via rounded rectangle
    # number
    add_text_box(slide, left + Pt(15), top + Pt(12), width - Pt(30), Pt(40),
                 str(number), font_name=TITLE_FONT, font_size=28, bold=True, color=num_color,
                 alignment=PP_ALIGN.CENTER)
    # label
    add_text_box(slide, left + Pt(10), top + Pt(55), width - Pt(20), Pt(30),
                 label, font_name=BODY_FONT, font_size=11, bold=False, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)
    return r


def add_section_divider(prs, title_text, subtitle_text=""):
    """Add a section divider slide (dark background, big serif title)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BG)
    # decorative line top
    add_accent_line(slide, Inches(0), Inches(0), Inches(13.33), ACCENT_GOLD, Pt(4))
    # gold accent bar left
    add_shape(slide, Inches(0.8), Inches(2.2), Pt(6), Inches(2.5), fill_color=ACCENT_GOLD)
    # title
    add_text_box(slide, Inches(1.3), Inches(2.3), Inches(10), Inches(1.2),
                 title_text, font_name=TITLE_FONT, font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT, line_spacing=1.3)
    if subtitle_text:
        add_text_box(slide, Inches(1.3), Inches(3.6), Inches(10), Inches(0.8),
                     subtitle_text, font_name=BODY_FONT, font_size=18, bold=False, color=MID_GRAY,
                     alignment=PP_ALIGN.LEFT)
    # bottom decorative line
    add_accent_line(slide, Inches(0.8), Inches(5.2), Inches(2), ACCENT_GOLD, Pt(2))
    # page number
    add_text_box(slide, Inches(12), Inches(6.8), Inches(1), Inches(0.3),
                 f"{len(prs.slides)}", font_name=TITLE_FONT, font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)
    return slide


# ── CREATE PRESENTATION ────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 - COVER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

# Large decorative top band
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.15), fill_color=ACCENT_GOLD)
# Left gold vertical bar
add_shape(slide, Inches(0.8), Inches(1.8), Pt(8), Inches(3.8), fill_color=ACCENT_GOLD)
# Main title
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.5),
             "AI Agent 驱动", font_name=TITLE_FONT, font_size=48, bold=True, color=WHITE,
             alignment=PP_ALIGN.LEFT, line_spacing=1.1)
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1.0),
             "商业洞察与云方案实践", font_name=TITLE_FONT, font_size=36, bold=True, color=ACCENT_GOLD,
             alignment=PP_ALIGN.LEFT, line_spacing=1.1)
# Subtitle
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(8), Inches(0.6),
             "智云赢商·AI Agent 创新实战大比武", font_name=BODY_FONT, font_size=20, bold=False, color=LIGHT_GRAY,
             alignment=PP_ALIGN.LEFT)
# Speaker & date
add_accent_line(slide, Inches(1.5), Inches(5.2), Inches(2.5), ACCENT_GOLD, Pt(2))
add_text_box(slide, Inches(1.5), Inches(5.4), Inches(5), Inches(0.4),
             "唐潘 · 解决方案", font_name=BODY_FONT, font_size=16, bold=True, color=WHITE,
             alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1.5), Inches(5.9), Inches(5), Inches(0.4),
             "2026", font_name=TITLE_FONT, font_size=14, bold=False, color=MID_GRAY,
             alignment=PP_ALIGN.LEFT)
# Bottom decorative
add_shape(slide, Inches(0), Inches(7.35), Inches(13.33), Inches(0.15), fill_color=ACCENT_GOLD)
# Decorative right side circles
for i, (y, sz) in enumerate([(1.5, 0.3), (2.3, 0.15), (6.0, 0.2)]):
    c = add_shape(slide, Inches(11.5), Inches(y), Inches(sz * 2.54), Inches(sz * 2.54),
                  fill_color=None, shape_type=MSO_SHAPE.OVAL)
    c.fill.solid()
    c.fill.fore_color.rgb = ACCENT_BLUE if i % 2 == 0 else ACCENT_GOLD
    c.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 - AGENDA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(5), Inches(0.7),
             "议程概览", font_name=TITLE_FONT, font_size=36, bold=True, color=DARK_TEXT,
             alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(1.5), ACCENT_GOLD, Pt(3))

agenda_items = [
    ("01", "背景与挑战", "从信息过载到智能决策"),
    ("02", "项目全景·双轨并行", "InsightPro + SAC 解决方案实践"),
    ("03", "InsightPro·商业洞察平台", "企业级AI商业洞察平台"),
    ("04", "SAC·解决方案实践", "基于RFS/OpenTofu的自动化方案"),
    ("05", "华为云产品深度使用", "MaaS · RFS · Headroom"),
    ("06", "创新性与业务价值总结", "双AI引擎 · 上下文压缩 · 6大场景"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    x = 0.8 + (i % 3) * 4.1
    y = 1.8 + (i // 3) * 2.6
    # card
    add_shape(slide, Inches(x), Inches(y), Inches(3.7), Inches(2.1),
              fill_color=CARD_BG, line_color=CARD_BORDER)
    # number circle
    c = add_shape(slide, Inches(x + 0.2), Inches(y + 0.2), Inches(0.5), Inches(0.5),
                  fill_color=ACCENT_GOLD, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.25), Inches(0.5), Inches(0.4),
                 num, font_name=TITLE_FONT, font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    # title
    add_text_box(slide, Inches(x + 0.9), Inches(y + 0.25), Inches(2.6), Inches(0.5),
                 title, font_name=BODY_FONT, font_size=16, bold=True, color=DARK_TEXT,
                 alignment=PP_ALIGN.LEFT)
    # description
    add_text_box(slide, Inches(x + 0.9), Inches(y + 0.8), Inches(2.6), Inches(1.0),
                 desc, font_name=BODY_FONT, font_size=11, bold=False, color=MID_GRAY,
                 alignment=PP_ALIGN.LEFT, line_spacing=1.4)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 - ACT 1 DIVIDER
# ═══════════════════════════════════════════════════════════════
add_section_divider(prs,
    "AI Agent 时代",
    "从信息过载到智能决策，从重复劳动到自动化开发")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 - DUAL CHALLENGES (left/right)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
             "双重挑战", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

# Left card - 情报获取之困
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2),
          fill_color=DARK_BG)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.6),
             "情报获取之困", font_name=TITLE_FONT, font_size=26, bold=True, color=ACCENT_GOLD)
add_accent_line(slide, Inches(1.1), Inches(2.4), Inches(1), ACCENT_BLUE, Pt(2))

left_items = [
    "信息散落各处 — 碎片化，海量多源",
    "人工搜集时效性差 — 滞后于市场变化",
    "知识无法沉淀 — 洞察随人员流动流失",
    "格式不统一 — PDF/网页/图片/公众号",
    "重复劳动 — 60%时间花在数据整理上",
]
for i, item in enumerate(left_items):
    add_text_box(slide, Inches(1.3), Inches(2.7 + i * 0.55), Inches(4.5), Inches(0.5),
                 f"▸ {item}", font_name=BODY_FONT, font_size=13, color=LIGHT_GRAY,
                 line_spacing=1.3)

# Right card - 开发部署之痛
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.2),
          fill_color=DARK_BG2)
add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.6),
             "开发部署之痛", font_name=TITLE_FONT, font_size=26, bold=True, color=ACCENT_CYAN)
add_accent_line(slide, Inches(7.3), Inches(2.4), Inches(1), ACCENT_CYAN, Pt(2))

right_items = [
    "方案周期漫长 — 从设计到部署需数周",
    "基础设施重复配置 — 每次手动搭建环境",
    "Token成本高 — 大代码上下文消耗惊人",
    "环境不一致 — 开发/测试/生产割裂",
    "部署复杂 — 多步骤手动操作易出错",
]
for i, item in enumerate(right_items):
    add_text_box(slide, Inches(7.5), Inches(2.7 + i * 0.55), Inches(4.5), Inches(0.5),
                 f"▸ {item}", font_name=BODY_FONT, font_size=13, color=LIGHT_GRAY,
                 line_spacing=1.3)

# Center connector arrow
add_shape(slide, Inches(6.45), Inches(3.7), Inches(0.45), Inches(0.45),
          fill_color=ACCENT_GOLD, shape_type=MSO_SHAPE.OVAL)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 - PROJECT PANORAMA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.7),
             "项目全景 · 双轨并行", font_name=TITLE_FONT, font_size=36, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(1.5), ACCENT_GOLD, Pt(3))

# Project A card
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5),
          fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5),
             "Project A", font_name=TITLE_FONT, font_size=14, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(0.6),
             "InsightPro", font_name=TITLE_FONT, font_size=28, bold=True, color=WHITE)
add_text_box(slide, Inches(1.1), Inches(2.8), Inches(5), Inches(0.5),
             "企业级AI商业洞察平台", font_name=BODY_FONT, font_size=16, bold=False, color=MID_GRAY)
add_accent_line(slide, Inches(1.1), Inches(3.4), Inches(2), ACCENT_BLUE, Pt(2))

proj_a_lines = [
    {"text": "技术栈", "size": 13, "bold": True, "color": ACCENT_GOLD},
    {"text": "Next.js 14 · FastAPI · DeepSeek API · MaaS", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "", "size": 6, "bold": False, "color": LIGHT_GRAY},
    {"text": "核心能力", "size": 13, "bold": True, "color": ACCENT_GOLD},
    {"text": "商业情报采集 · AI分析 · 可视化展示", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "", "size": 6, "bold": False, "color": LIGHT_GRAY},
    {"text": "部署状态", "size": 13, "bold": True, "color": GREEN_ACCENT},
    {"text": "已上线 94.74.90.21:4000", "size": 12, "bold": False, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(1.1), Inches(3.6), Inches(5), Inches(2.5), proj_a_lines,
               line_spacing=1.4)

# Project B card
add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.5),
          fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=ACCENT_CYAN)
add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.5),
             "Project B", font_name=TITLE_FONT, font_size=14, bold=True, color=ACCENT_CYAN)
add_text_box(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(0.6),
             "SAC 解决方案实践", font_name=TITLE_FONT, font_size=28, bold=True, color=WHITE)
add_text_box(slide, Inches(7.3), Inches(2.8), Inches(5), Inches(0.5),
             "基于RFS/OpenTofu的自动化方案", font_name=BODY_FONT, font_size=16, bold=False, color=MID_GRAY)
add_accent_line(slide, Inches(7.3), Inches(3.4), Inches(2), ACCENT_CYAN, Pt(2))

proj_b_lines = [
    {"text": "技术栈", "size": 13, "bold": True, "color": ACCENT_GOLD},
    {"text": "RFS · OpenTofu · Headroom · MaaS", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "", "size": 6, "bold": False, "color": LIGHT_GRAY},
    {"text": "核心能力", "size": 13, "bold": True, "color": ACCENT_GOLD},
    {"text": "一键部署 · Token压缩 · 多云编排", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "", "size": 6, "bold": False, "color": LIGHT_GRAY},
    {"text": "覆盖范围", "size": 13, "bold": True, "color": GREEN_ACCENT},
    {"text": "CN + HK 双区域", "size": 12, "bold": False, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(7.3), Inches(3.6), Inches(5), Inches(2.5), proj_b_lines,
               line_spacing=1.4)

# Bottom stats bar
add_shape(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
          fill_color=RGBColor(0x1E, 0x1E, 0x3A))
stats = [("2", "项目"), ("3", "华为云服务\nMaaS·RFS·Headroom"), ("100%", "开源")]
for i, (num, label) in enumerate(stats):
    x = 1.0 + i * 4.0
    add_text_box(slide, Inches(x), Inches(6.35), Inches(1.0), Inches(0.6),
                 num, font_name=TITLE_FONT, font_size=24, bold=True, color=ACCENT_GOLD,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.9), Inches(6.35), Inches(2.8), Inches(0.6),
                 label, font_name=BODY_FONT, font_size=11, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.LEFT)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 - ACT 2 DIVIDER: InsightPro
# ═══════════════════════════════════════════════════════════════
add_section_divider(prs, "InsightPro", "让数据说话")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 - InsightPro Business Value (6 data cards)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "InsightPro 业务价值", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

data_cards = [
    ("6", "大行业覆盖", ACCENT_BLUE),
    ("5", "大厂商对比", PURPLE_ACCENT),
    ("12", "个功能模块", GREEN_ACCENT),
    ("24h", "实时监控", RED_ACCENT),
    ("16", "个API端点", ACCENT_CYAN),
    ("已部署上线", "94.74.90.21:4000", ACCENT_GOLD),
]

for i, (num, label, clr) in enumerate(data_cards):
    x = 0.6 + (i % 3) * 4.2
    y = 1.5 + (i // 3) * 2.8
    add_data_card(slide, Inches(x), Inches(y), Inches(3.8), Inches(2.3),
                  num, label, num_color=clr)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 - THREE-LAYER TECH ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "三层技术架构", font_name=TITLE_FONT, font_size=34, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

layers = [
    ("Layer 1", "数据采集层", ACCENT_CYAN, RGBColor(0x07, 0x2F, 0x3F),
     ["百度热搜 · GitHub Trending", "竞品数据爬取 · RSS订阅", "缓存降级策略 (SQLite)", "APScheduler定时采集"]),
    ("Layer 2", "AI分析层（华为云MaaS）", ACCENT_BLUE, RGBColor(0x0C, 0x24, 0x4A),
     ["DeepSeek API 结构化分析", "ModelArts MaaS 流式对话", "APScheduler任务编排", "优雅降级机制保障可用性"]),
    ("Layer 3", "展示交互层", ACCENT_GOLD, RGBColor(0x3A, 0x2A, 0x0A),
     ["Next.js 14 前端框架", "FastAPI 后端服务", "Recharts 数据可视化", "Tailwind CSS 响应式设计"]),
]

for i, (layer_num, layer_name, accent, layer_bg, items) in enumerate(layers):
    y = 1.4 + i * 1.9
    # layer card
    add_shape(slide, Inches(0.8), Inches(y), Inches(11.7), Inches(1.6),
              fill_color=layer_bg, line_color=accent)
    # layer number badge
    badge = add_shape(slide, Inches(0.9), Inches(y + 0.1), Inches(0.6), Inches(0.35),
                      fill_color=accent, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.9), Inches(y + 0.12), Inches(0.6), Inches(0.3),
                 layer_num, font_name=TITLE_FONT, font_size=11, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    # layer name
    add_text_box(slide, Inches(1.7), Inches(y + 0.1), Inches(4), Inches(0.4),
                 layer_name, font_name=BODY_FONT, font_size=16, bold=True, color=WHITE)
    # items
    for j, item in enumerate(items):
        ix = 1.7 + (j % 2) * 5.5
        iy = y + 0.6 + (j // 2) * 0.5
        add_text_box(slide, Inches(ix), Inches(iy), Inches(5), Inches(0.4),
                     f"▸ {item}", font_name=BODY_FONT, font_size=12, color=LIGHT_GRAY)

# Arrow connectors between layers
for i in range(2):
    y = 3.0 + i * 1.9
    add_shape(slide, Inches(6.5), Inches(y), Inches(0.3), Inches(0.3),
              fill_color=ACCENT_GOLD, shape_type=MSO_SHAPE.DOWN_ARROW)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 - 12 FUNCTIONAL MODULES (6-grid)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(6), Inches(0.6),
             "12 大功能模块", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.2), ACCENT_GOLD, Pt(3))

modules = [
    ("今日洞察", "AI每日情报摘要", ACCENT_BLUE),
    ("行业全景", "多维度行业分析", GREEN_ACCENT),
    ("案例库", "标杆案例研究", PURPLE_ACCENT),
    ("热点追踪", "实时趋势监控", RED_ACCENT),
    ("友商洞察", "竞品动态分析", ACCENT_CYAN),
    ("政策法规", "政策变化追踪", ACCENT_GOLD),
    ("商业快讯", "行业新闻聚合", ACCENT_BLUE),
    ("增长机会", "市场机会识别", GREEN_ACCENT),
    ("市场情报", "市场趋势分析", PURPLE_ACCENT),
    ("需求挖掘", "客户需求洞察", RED_ACCENT),
    ("招标信息", "招投标情报", ACCENT_CYAN),
    ("数据大屏", "可视化仪表盘", ACCENT_GOLD),
]

for i, (name, desc, clr) in enumerate(modules):
    x = 0.5 + (i % 4) * 3.15
    y = 1.3 + (i // 4) * 2.0
    # card
    add_shape(slide, Inches(x), Inches(y), Inches(2.9), Inches(1.7),
              fill_color=CARD_BG, line_color=CARD_BORDER)
    # top accent bar
    add_shape(slide, Inches(x), Inches(y), Inches(2.9), Pt(4), fill_color=clr)
    # number
    add_text_box(slide, Inches(x + 0.15), Inches(y + 0.15), Inches(0.4), Inches(0.4),
                 f"{i+1:02d}", font_name=TITLE_FONT, font_size=16, bold=True, color=clr)
    # name
    add_text_box(slide, Inches(x + 0.6), Inches(y + 0.15), Inches(2.1), Inches(0.4),
                 name, font_name=BODY_FONT, font_size=15, bold=True, color=DARK_TEXT)
    # desc
    add_text_box(slide, Inches(x + 0.15), Inches(y + 0.7), Inches(2.6), Inches(0.8),
                 desc, font_name=BODY_FONT, font_size=11, color=MID_GRAY, line_spacing=1.3)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 - SIX-DIMENSION COMPETITIVE LANDSCAPE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "六维竞争场景", font_name=TITLE_FONT, font_size=34, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
             "华为云 vs AWS · Azure · 阿里云 · 腾讯云 · 火山云", font_name=BODY_FONT, font_size=13, color=MID_GRAY)

scenarios = [
    ("中资企业出海", "全球化部署与合规", ACCENT_BLUE),
    ("AI大模型平台", "ModelArts MaaS能力", PURPLE_ACCENT),
    ("开发者生态", "工具链与社区支持", GREEN_ACCENT),
    ("企业SaaS", "企业级应用场景", RED_ACCENT),
    ("腰部客户性价比", "成本与性能平衡", ACCENT_GOLD),
    ("传统行业数字化", "行业数字化转型", ACCENT_CYAN),
]

for i, (name, desc, clr) in enumerate(scenarios):
    x = 0.6 + (i % 3) * 4.2
    y = 1.8 + (i // 3) * 2.6
    add_shape(slide, Inches(x), Inches(y), Inches(3.8), Inches(2.2),
              fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=clr)
    # Icon circle
    c = add_shape(slide, Inches(x + 1.4), Inches(y + 0.2), Inches(0.5), Inches(0.5),
                  fill_color=clr, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(x + 0.3), Inches(y + 0.85), Inches(3.2), Inches(0.4),
                 name, font_name=BODY_FONT, font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.3), Inches(y + 1.35), Inches(3.2), Inches(0.6),
                 desc, font_name=BODY_FONT, font_size=12, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 - DeepSeek + MaaS Dual Engine
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "DeepSeek + MaaS 双引擎", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

# Engine A
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
          fill_color=DARK_BG, line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5),
             "Engine A", font_name=TITLE_FONT, font_size=14, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(0.6),
             "DeepSeek API", font_name=TITLE_FONT, font_size=26, bold=True, color=WHITE)
add_text_box(slide, Inches(1.1), Inches(2.7), Inches(5), Inches(0.4),
             "结构化分析引擎", font_name=BODY_FONT, font_size=14, color=MID_GRAY)
add_accent_line(slide, Inches(1.1), Inches(3.2), Inches(2), ACCENT_BLUE, Pt(2))

engine_a_lines = [
    {"text": "▸ 深度商业情报结构化分析", "size": 13, "bold": False, "color": LIGHT_GRAY},
    {"text": "▸ 多维度数据交叉对比", "size": 13, "bold": False, "color": LIGHT_GRAY},
    {"text": "▸ 自动化报告生成", "size": 13, "bold": False, "color": LIGHT_GRAY},
    {"text": "▸ 批量处理与企业级吞吐", "size": 13, "bold": False, "color": LIGHT_GRAY},
    {"text": "▸ 降级机制保障高可用", "size": 13, "bold": False, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(1.1), Inches(3.5), Inches(5), Inches(2.8),
               engine_a_lines, line_spacing=1.6)

# Engine B
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.0),
          fill_color=DARK_BG2, line_color=ACCENT_GOLD)
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
             "Engine B", font_name=TITLE_FONT, font_size=14, bold=True, color=ACCENT_GOLD)
add_text_box(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(0.6),
             "ModelArts MaaS", font_name=TITLE_FONT, font_size=26, bold=True, color=WHITE)
add_text_box(slide, Inches(7.3), Inches(2.7), Inches(5), Inches(0.4),
             "流式对话引擎 (deepseek-v3.2)", font_name=BODY_FONT, font_size=14, color=MID_GRAY)
add_accent_line(slide, Inches(7.3), Inches(3.2), Inches(2), ACCENT_GOLD, Pt(2))

engine_b_lines = [
    {"text": "▸ 实时流式对话交互", "size": 13, "bold": False, "color": LIGHT_GRAY},
    {"text": "▸ deepseek-v3.2 模型推理", "size": 13, "bold": False, "color": LIGHT_GRAY},
    {"text": "▸ Huawei Cloud原生托管", "size": 13, "bold": False, "color": LIGHT_GRAY},
    {"text": "▸ 弹性扩缩容", "size": 13, "bold": False, "color": LIGHT_GRAY},
    {"text": "▸ 企业级安全与合规", "size": 13, "bold": False, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(7.3), Inches(3.5), Inches(5), Inches(2.8),
               engine_b_lines, line_spacing=1.6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 - PRODUCT FEEDBACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "InsightPro 产品反馈", font_name=TITLE_FONT, font_size=34, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

# Big data highlights
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.0),
          fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=GREEN_ACCENT)
add_text_box(slide, Inches(1.1), Inches(1.6), Inches(2), Inches(0.6),
             "80%", font_name=TITLE_FONT, font_size=40, bold=True, color=GREEN_ACCENT)
add_text_box(slide, Inches(3.0), Inches(1.7), Inches(3), Inches(0.8),
             "信息聚合效率提升", font_name=BODY_FONT, font_size=16, bold=True, color=WHITE)
add_text_box(slide, Inches(1.1), Inches(2.4), Inches(5), Inches(0.5),
             "从人工搜集到AI自动聚合，效率飞跃", font_name=BODY_FONT, font_size=12, color=MID_GRAY)

add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.0),
          fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=ACCENT_BLUE)
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(2.8), Inches(0.5),
             "竞争分析深度结构化", font_name=BODY_FONT, font_size=18, bold=True, color=ACCENT_BLUE)
add_text_box(slide, Inches(7.3), Inches(2.3), Inches(4.5), Inches(0.8),
             "五大厂商、六大场景、多维度对比分析，从数据到洞察一步到位",
             font_name=BODY_FONT, font_size=12, color=LIGHT_GRAY, line_spacing=1.4)

# Future direction
add_shape(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(3.0),
          fill_color=RGBColor(0x1E, 0x1E, 0x3A))
add_text_box(slide, Inches(1.1), Inches(3.9), Inches(5), Inches(0.5),
             "持续迭代方向", font_name=TITLE_FONT, font_size=20, bold=True, color=ACCENT_GOLD)
add_accent_line(slide, Inches(1.1), Inches(4.4), Inches(1.5), ACCENT_GOLD, Pt(2))

future_items = [
    ("Supabase", "持久化存储", ACCENT_GREEN := GREEN_ACCENT),
    ("用户认证", "RBAC权限管理", ACCENT_BLUE),
    ("全局搜索", "全文检索能力", PURPLE_ACCENT),
    ("通知系统", "实时消息推送", ACCENT_CYAN),
    ("PDF导出", "报告生成", RED_ACCENT),
    ("移动端适配", "响应式Web", ACCENT_GOLD),
]
for i, (name, desc, clr) in enumerate(future_items):
    x = 1.0 + (i % 3) * 4.0
    y = 4.7 + (i // 3) * 0.9
    # bullet + name
    add_shape(slide, Inches(x), Inches(y + 0.05), Pt(6), Pt(6), fill_color=clr, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(x + 0.2), Inches(y), Inches(1.2), Inches(0.3),
                 name, font_name=BODY_FONT, font_size=13, bold=True, color=WHITE)
    add_text_box(slide, Inches(x + 1.4), Inches(y), Inches(2.3), Inches(0.3),
                 desc, font_name=BODY_FONT, font_size=12, color=MID_GRAY)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 - ACT 3 DIVIDER: SAC
# ═══════════════════════════════════════════════════════════════
add_section_divider(prs, "SAC 解决方案实践", "Solution Practices")

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 - SIX SOLUTIONS PANORAMA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(6), Inches(0.6),
             "六大解决方案全景", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.2), ACCENT_GOLD, Pt(3))

solutions = [
    ("Headroom\nOpenCode", "CN + HK", ACCENT_BLUE),
    ("LiteLLM", "CN + HK", PURPLE_ACCENT),
    ("OpenHands", "CN", GREEN_ACCENT),
    ("Supabase", "CN", RED_ACCENT),
    ("CodeWhale", "CN", ACCENT_CYAN),
    ("AiToEarn", "CN + HK", ACCENT_GOLD),
]

for i, (name, region, clr) in enumerate(solutions):
    x = 0.5 + (i % 3) * 4.2
    y = 1.3 + (i // 3) * 2.5
    add_shape(slide, Inches(x), Inches(y), Inches(3.8), Inches(2.1),
              fill_color=CARD_BG, line_color=clr)
    # top accent
    add_shape(slide, Inches(x), Inches(y), Inches(3.8), Pt(5), fill_color=clr)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.2), Inches(3.4), Inches(0.8),
                 name, font_name=BODY_FONT, font_size=18, bold=True, color=DARK_TEXT,
                 line_spacing=1.2)
    # region badge
    badge = add_shape(slide, Inches(x + 0.2), Inches(y + 1.2), Inches(1.0), Inches(0.3),
                      fill_color=clr, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 1.2), Inches(1.0), Inches(0.3),
                 region, font_name=BODY_FONT, font_size=10, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 1.6), Inches(3.4), Inches(0.4),
                 "统一目录结构", font_name=BODY_FONT, font_size=11, color=MID_GRAY)

# Bottom unified structure note
add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.7),
          fill_color=DARK_BG)
add_text_box(slide, Inches(1.1), Inches(6.25), Inches(11), Inches(0.5),
             "统一目录结构：docs · scripts · terraform",
             font_name=BODY_FONT, font_size=13, bold=True, color=ACCENT_GOLD,
             alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 - Headroom OpenCode Core Value
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "Headroom OpenCode 核心方案", font_name=TITLE_FONT, font_size=34, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

values = [
    ("60-95%", "Token压缩", "大幅降低AI调用成本", ACCENT_BLUE),
    ("零费用", "许可证费用", "开源MIT协议，自由使用", GREEN_ACCENT),
    ("一键部署", "Flexus X", "自动化部署，3分钟上线", ACCENT_GOLD),
    ("多模型", "DeepSeek+Claude", "灵活切换，最优选择", PURPLE_ACCENT),
]

for i, (num, title, desc, clr) in enumerate(values):
    x = 0.6 + (i % 2) * 6.3
    y = 1.5 + (i // 2) * 2.7
    add_shape(slide, Inches(x), Inches(y), Inches(5.8), Inches(2.3),
              fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=clr)
    add_text_box(slide, Inches(x + 0.3), Inches(y + 0.3), Inches(2.5), Inches(0.8),
                 num, font_name=TITLE_FONT, font_size=38, bold=True, color=clr)
    add_text_box(slide, Inches(x + 2.8), Inches(y + 0.3), Inches(2.8), Inches(0.5),
                 title, font_name=BODY_FONT, font_size=20, bold=True, color=WHITE)
    add_text_box(slide, Inches(x + 2.8), Inches(y + 0.9), Inches(2.8), Inches(0.5),
                 desc, font_name=BODY_FONT, font_size=12, color=MID_GRAY)
    add_accent_line(slide, Inches(x + 0.3), Inches(y + 1.5), Inches(5.2), clr, Pt(1))

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 - Headroom Tech Architecture
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "Headroom 技术架构", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.2), ACCENT_GOLD, Pt(3))

# Architecture flow - horizontal
boxes = [
    ("开发者\nOpenCode CLI", ACCENT_BLUE),
    ("→", ACCENT_GOLD),
    ("Headroom\nProxy (:8787)", ACCENT_CYAN),
    ("→", ACCENT_GOLD),
    ("华为云\nMaaS API", GREEN_ACCENT),
    ("→", ACCENT_GOLD),
    ("LLM\n推理", PURPLE_ACCENT),
]

for i, (label, clr) in enumerate(boxes):
    if label == "→":
        add_text_box(slide, Inches(1.2 + i * 1.8), Inches(1.3), Inches(0.5), Inches(0.8),
                     label, font_name=TITLE_FONT, font_size=24, bold=True, color=ACCENT_GOLD,
                     alignment=PP_ALIGN.CENTER)
    else:
        x = 0.8 + i * 1.8
        add_shape(slide, Inches(x), Inches(1.2), Inches(1.5), Inches(1.0),
                  fill_color=DARK_BG, line_color=clr)
        add_text_box(slide, Inches(x), Inches(1.25), Inches(1.5), Inches(0.9),
                     label, font_name=BODY_FONT, font_size=11, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Resources section
add_shape(slide, Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.2),
          fill_color=DARK_BG)
add_text_box(slide, Inches(1.1), Inches(2.7), Inches(5), Inches(0.5),
             "部署资源清单", font_name=TITLE_FONT, font_size=20, bold=True, color=ACCENT_GOLD)

resources = [
    ("Flexus X ECS", "x1.2u.4g", "弹性计算资源", ACCENT_BLUE),
    ("VPC", "172.16.0.0/16", "虚拟私有网络", ACCENT_CYAN),
    ("安全组", "22 + 8787", "端口策略配置", GREEN_ACCENT),
    ("弹性IP", "公网访问", "动态/静态绑定", ACCENT_GOLD),
]

for i, (name, spec, desc, clr) in enumerate(resources):
    x = 1.0 + (i % 2) * 5.8
    y = 3.5 + (i // 2) * 1.4
    add_shape(slide, Inches(x), Inches(y), Inches(5.3), Inches(1.1),
              fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=clr)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.1), Inches(1.5), Inches(0.4),
                 name, font_name=BODY_FONT, font_size=15, bold=True, color=WHITE)
    add_text_box(slide, Inches(x + 1.8), Inches(y + 0.1), Inches(2.0), Inches(0.4),
                 spec, font_name=BODY_FONT, font_size=13, bold=True, color=clr)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.55), Inches(4.5), Inches(0.4),
                 desc, font_name=BODY_FONT, font_size=11, color=MID_GRAY)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 - RFS Resource Orchestration
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "RFS 资源编排", font_name=TITLE_FONT, font_size=34, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

# Template structure - left
add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.3),
          fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(1.5), Inches(5), Inches(0.5),
             "模板结构", font_name=TITLE_FONT, font_size=18, bold=True, color=ACCENT_BLUE)

template_items = [
    "Variables — 输入参数定义",
    "Locals — 本地变量计算",
    "Data Sources — 数据源查询",
    "VPC / Subnet — 网络配置",
    "SecGroup — 安全组策略",
    "EIP — 弹性IP分配",
    "ECS — 云服务器实例",
    "user_data — 启动脚本注入",
]
for i, item in enumerate(template_items):
    add_text_box(slide, Inches(1.3), Inches(2.2 + i * 0.45), Inches(4.5), Inches(0.4),
                 f"▸ {item}", font_name=BODY_FONT, font_size=12, color=LIGHT_GRAY)

# Right side - principles + flow
add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(2.0),
          fill_color=RGBColor(0x1E, 0x1E, 0x3A), line_color=ACCENT_GOLD)
add_text_box(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.4),
             "设计原则", font_name=TITLE_FONT, font_size=16, bold=True, color=ACCENT_GOLD)
principles = [
    "模板-脚本分离 — 关注点分离",
    "模块化设计 — 可复用组件",
    "参数化配置 — 灵活适配多环境",
]
for i, p in enumerate(principles):
    add_text_box(slide, Inches(7.3), Inches(2.0 + i * 0.4), Inches(4.5), Inches(0.35),
                 f"✦ {p}", font_name=BODY_FONT, font_size=12, color=LIGHT_GRAY)

# Flow
add_shape(slide, Inches(7.0), Inches(3.7), Inches(5.5), Inches(3.0),
          fill_color=RGBColor(0x1E, 0x1E, 0x3A), line_color=GREEN_ACCENT)
add_text_box(slide, Inches(7.3), Inches(3.8), Inches(5), Inches(0.4),
             "RFS 一键部署流程", font_name=TITLE_FONT, font_size=16, bold=True, color=GREEN_ACCENT)

flow_steps = [
    ("1", "打开RFS链接", ACCENT_BLUE),
    ("2", "填参数", ACCENT_CYAN),
    ("3", "一键部署", GREEN_ACCENT),
    ("4", "验证", ACCENT_GOLD),
]
for i, (num, step, clr) in enumerate(flow_steps):
    x = 7.3 + i * 1.3
    y = 4.4
    c = add_shape(slide, Inches(x + 0.15), Inches(y), Inches(0.4), Inches(0.4),
                  fill_color=clr, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(x + 0.15), Inches(y + 0.02), Inches(0.4), Inches(0.35),
                 num, font_name=TITLE_FONT, font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x), Inches(y + 0.5), Inches(1.2), Inches(0.3),
                 step, font_name=BODY_FONT, font_size=11, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, Inches(x + 0.9), Inches(y + 0.05), Inches(0.4), Inches(0.3),
                     "→", font_name=TITLE_FONT, font_size=16, bold=True, color=ACCENT_GOLD)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 - MULTI-SOLUTION MATRIX
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(6), Inches(0.6),
             "多方案矩阵", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.2), ACCENT_GOLD, Pt(3))

# Table-like layout
headers = ["方案名称", "区域", "核心用途", "部署方式", "状态"]
rows_data = [
    ["Headroom OpenCode", "CN / HK", "Token压缩代理", "RFS / 一键部署", "✓ 已上线"],
    ["LiteLLM", "CN / HK", "多模型代理网关", "RFS / 一键部署", "✓ 已上线"],
    ["OpenHands", "CN", "AI编程助手", "RFS / 一键部署", "✓ 已上线"],
    ["Supabase", "CN", "后端即服务", "RFS / 一键部署", "✓ 已上线"],
    ["CodeWhale", "CN", "代码分析平台", "RFS / 一键部署", "✓ 已上线"],
    ["AiToEarn", "CN / HK", "AI变现平台", "RFS / 一键部署", "✓ 已上线"],
]

# Header row
col_widths = [2.2, 1.5, 3.0, 2.5, 1.5]
x_start = 0.8
for j, (hdr, w) in enumerate(zip(headers, col_widths)):
    x = x_start + sum(col_widths[:j])
    add_shape(slide, Inches(x), Inches(1.2), Inches(w), Inches(0.5),
              fill_color=DARK_BG)
    add_text_box(slide, Inches(x), Inches(1.22), Inches(w), Inches(0.45),
                 hdr, font_name=BODY_FONT, font_size=13, bold=True, color=ACCENT_GOLD,
                 alignment=PP_ALIGN.CENTER)

# Data rows
for i, row in enumerate(rows_data):
    bg = RGBColor(0x16, 0x1B, 0x33) if i % 2 == 0 else RGBColor(0x1E, 0x1E, 0x3A)
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        x = x_start + sum(col_widths[:j])
        add_shape(slide, Inches(x), Inches(1.7 + i * 0.55), Inches(w), Inches(0.55),
                  fill_color=bg, line_color=RGBColor(0x2A, 0x2A, 0x4A))
        clr = GREEN_ACCENT if cell.startswith("✓") else WHITE
        add_text_box(slide, Inches(x + 0.1), Inches(1.75 + i * 0.55), Inches(w - 0.2), Inches(0.45),
                     cell, font_name=BODY_FONT, font_size=12, bold=(j == 0), color=clr,
                     alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 - SAC PRODUCT FEEDBACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "SAC 产品反馈", font_name=TITLE_FONT, font_size=34, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

feedback_cards = [
    ("60-80%", "日常开发", "代码编写、调试、重构场景下\nToken消耗显著降低", ACCENT_BLUE),
    ("85-92%", "大代码库探索", "大型项目中快速定位代码\n上下文压缩效果显著", ACCENT_CYAN),
    ("90-95%", "日志调试分析", "海量日志分析场景\n压缩比最高，效果最突出", GREEN_ACCENT),
]

for i, (pct, title, desc, clr) in enumerate(feedback_cards):
    x = 0.6 + i * 4.2
    add_shape(slide, Inches(x), Inches(1.4), Inches(3.8), Inches(4.5),
              fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=clr)
    # Top percentage
    add_shape(slide, Inches(x), Inches(1.4), Inches(3.8), Inches(1.5),
              fill_color=RGBColor(0x0C, 0x12, 0x28))
    add_text_box(slide, Inches(x + 0.2), Inches(1.5), Inches(3.4), Inches(0.8),
                 pct, font_name=TITLE_FONT, font_size=42, bold=True, color=clr,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.2), Inches(2.2), Inches(3.4), Inches(0.4),
                 title, font_name=BODY_FONT, font_size=18, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, Inches(x + 0.5), Inches(2.7), Inches(2.8), clr, Pt(1))
    # Description
    add_text_box(slide, Inches(x + 0.3), Inches(3.0), Inches(3.2), Inches(1.5),
                 desc, font_name=BODY_FONT, font_size=12, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 - ACT 4 DIVIDER: Huawei Cloud Deep Integration
# ═══════════════════════════════════════════════════════════════
add_section_divider(prs, "华为云产品深度集成", "MaaS · RFS · Headroom 三位一体")

# ═══════════════════════════════════════════════════════════════
# SLIDE 21 - MaaS Dual Project Coverage
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "MaaS 双项目全覆盖", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.2), ACCENT_GOLD, Pt(3))

# InsightPro MaaS
add_shape(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.2),
          fill_color=DARK_BG, line_color=ACCENT_CYAN)
add_text_box(slide, Inches(1.1), Inches(1.5), Inches(5), Inches(0.5),
             "InsightPro 中的 MaaS", font_name=TITLE_FONT, font_size=20, bold=True, color=ACCENT_CYAN)
add_accent_line(slide, Inches(1.1), Inches(2.0), Inches(2), ACCENT_CYAN, Pt(2))

ip_maas = [
    {"text": "流式对话引擎", "size": 16, "bold": True, "color": WHITE},
    {"text": "deepseek-v3.2 实时对话交互", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "华为云MaaS原生托管", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "", "size": 8, "bold": False, "color": LIGHT_GRAY},
    {"text": "结构化分析引擎", "size": 16, "bold": True, "color": WHITE},
    {"text": "DeepSeek API批量分析", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "缓存+降级双保障策略", "size": 12, "bold": False, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(4),
               ip_maas, line_spacing=1.4)

# SAC MaaS
add_shape(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(5.2),
          fill_color=DARK_BG2, line_color=ACCENT_GOLD)
add_text_box(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.5),
             "SAC 中的 MaaS", font_name=TITLE_FONT, font_size=20, bold=True, color=ACCENT_GOLD)
add_accent_line(slide, Inches(7.3), Inches(2.0), Inches(2), ACCENT_GOLD, Pt(2))

sac_maas = [
    {"text": "AI编程推理后端", "size": 16, "bold": True, "color": WHITE},
    {"text": "Headroom Proxy 网关集成", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "多模型路由 (DeepSeek+Claude)", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "", "size": 8, "bold": False, "color": LIGHT_GRAY},
    {"text": "Token压缩代理", "size": 16, "bold": True, "color": WHITE},
    {"text": "AST级别上下文压缩", "size": 12, "bold": False, "color": LIGHT_GRAY},
    {"text": "可逆压缩保障数据完整性", "size": 12, "bold": False, "color": LIGHT_GRAY},
]
add_multi_text(slide, Inches(7.3), Inches(2.2), Inches(5), Inches(4),
               sac_maas, line_spacing=1.4)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 22 - RFS 5-Step Dev Process
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "RFS 开发流程", font_name=TITLE_FONT, font_size=34, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

steps = [
    ("01", "需求分析", "明确部署目标\n确定资源规格\n选择区域配置", ACCENT_BLUE),
    ("02", "Terraform\n编写", "定义基础设施\n模板-脚本分离\n参数化配置", ACCENT_CYAN),
    ("03", "安装脚本", "编写user_data\n自动化安装\n环境初始化", GREEN_ACCENT),
    ("04", "OBS上传", "上传模板/脚本\n版本化管理\n权限配置", ACCENT_GOLD),
    ("05", "RFS部署\n验证", "一键部署堆栈\n状态监控\n功能验证", PURPLE_ACCENT),
]

for i, (num, title, desc, clr) in enumerate(steps):
    x = 0.5 + i * 2.55
    # Step card
    add_shape(slide, Inches(x), Inches(1.4), Inches(2.3), Inches(5.0),
              fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=clr)
    # Number circle
    c = add_shape(slide, Inches(x + 0.75), Inches(1.6), Inches(0.5), Inches(0.5),
                  fill_color=clr, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(x + 0.75), Inches(1.63), Inches(0.5), Inches(0.45),
                 num, font_name=TITLE_FONT, font_size=16, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, Inches(x + 0.1), Inches(2.3), Inches(2.1), Inches(0.6),
                 title, font_name=BODY_FONT, font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.2)
    # Accent line
    add_accent_line(slide, Inches(x + 0.4), Inches(3.0), Inches(1.5), clr, Pt(1))
    # Description
    add_text_box(slide, Inches(x + 0.15), Inches(3.2), Inches(2.0), Inches(2.0),
                 desc, font_name=BODY_FONT, font_size=11, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    # Arrow between steps
    if i < 4:
        add_text_box(slide, Inches(x + 2.2), Inches(3.5), Inches(0.4), Inches(0.4),
                     "›", font_name=TITLE_FONT, font_size=24, bold=True, color=ACCENT_GOLD,
                     alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 23 - Headroom Compression Capability Matrix
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "Headroom 压缩能力矩阵", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.2), ACCENT_GOLD, Pt(3))

# Matrix cards
matrix_items = [
    ("60-80%", "日常开发", "代码编写、调试、重构\n上下文智能压缩", ACCENT_BLUE),
    ("85-92%", "大代码库探索", "大型项目中快速定位\n代码结构优化压缩", PURPLE_ACCENT),
    ("70-85%", "架构设计", "架构文档与设计讨论\n聚焦关键信息", ACCENT_CYAN),
    ("90-95%", "日志分析调试", "海量日志秒级分析\n极致压缩比表现", GREEN_ACCENT),
]

for i, (pct, title, desc, clr) in enumerate(matrix_items):
    x = 0.5 + (i % 2) * 6.3
    y = 1.3 + (i // 2) * 2.8
    add_shape(slide, Inches(x), Inches(y), Inches(5.8), Inches(2.4),
              fill_color=CARD_BG, line_color=clr)
    # Left percentage
    add_shape(slide, Inches(x), Inches(y), Inches(2.2), Inches(2.4),
              fill_color=DARK_BG)
    add_text_box(slide, Inches(x + 0.1), Inches(y + 0.3), Inches(2.0), Inches(0.8),
                 pct, font_name=TITLE_FONT, font_size=36, bold=True, color=clr,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.1), Inches(y + 1.1), Inches(2.0), Inches(0.4),
                 "压缩率", font_name=BODY_FONT, font_size=12, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)
    # Right content
    add_text_box(slide, Inches(x + 2.5), Inches(y + 0.3), Inches(3.0), Inches(0.5),
                 title, font_name=BODY_FONT, font_size=18, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(x + 2.5), Inches(y + 0.9), Inches(3.0), Inches(1.0),
                 desc, font_name=BODY_FONT, font_size=12, color=MID_GRAY, line_spacing=1.4)

# Tech highlights bar
add_shape(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7),
          fill_color=DARK_BG)
highlights = "AST级别压缩 · 可逆压缩 · Prometheus监控 · 数据主权 · 多模型支持"
add_text_box(slide, Inches(1.1), Inches(6.35), Inches(11), Inches(0.5),
             f"技术亮点：{highlights}", font_name=BODY_FONT, font_size=13, bold=True, color=ACCENT_GOLD,
             alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 24 - ACT 5 DIVIDER: Innovation Drives Value
# ═══════════════════════════════════════════════════════════════
add_section_divider(prs, "创新驱动价值", "从技术突破到商业成果")

# ═══════════════════════════════════════════════════════════════
# SLIDE 25 - AI Agent Innovation (6 items)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "AI Agent 创新点", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.2), ACCENT_GOLD, Pt(3))

innovations = [
    ("双AI引擎协同", "DeepSeek结构化分析 + MaaS流式对话\n双引擎互补，覆盖全场景", ACCENT_BLUE),
    ("上下文压缩革命", "AST级别可逆压缩\nToken成本降低60-95%", PURPLE_ACCENT),
    ("解决方案实践", "RFS + OpenTofu一键部署\n从需求到上线流程标准化", GREEN_ACCENT),
    ("Agent辅助竞争分析", "AI自动采集对比五大厂商\n六大场景深度结构化", ACCENT_CYAN),
    ("自动化情报流水线", "APScheduler定时采集分析\n从数据到洞察全自动化", ACCENT_GOLD),
    ("Skill技能市场", "可复用AI Skill组件\n持续积累组织智慧", RED_ACCENT),
]

for i, (title, desc, clr) in enumerate(innovations):
    x = 0.5 + (i % 3) * 4.2
    y = 1.3 + (i // 3) * 2.8
    add_shape(slide, Inches(x), Inches(y), Inches(3.8), Inches(2.4),
              fill_color=CARD_BG, line_color=clr)
    # Top accent
    add_shape(slide, Inches(x), Inches(y), Inches(3.8), Pt(5), fill_color=clr)
    # Icon/number
    c = add_shape(slide, Inches(x + 0.15), Inches(y + 0.2), Inches(0.4), Inches(0.4),
                  fill_color=clr, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(x + 0.15), Inches(y + 0.23), Inches(0.4), Inches(0.35),
                 f"{i+1}", font_name=TITLE_FONT, font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, Inches(x + 0.7), Inches(y + 0.2), Inches(2.9), Inches(0.4),
                 title, font_name=BODY_FONT, font_size=16, bold=True, color=DARK_TEXT)
    # Description
    add_text_box(slide, Inches(x + 0.15), Inches(y + 0.8), Inches(3.5), Inches(1.3),
                 desc, font_name=BODY_FONT, font_size=12, color=MID_GRAY, line_spacing=1.5)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 26 - BUSINESS VALUE SUMMARY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
             "业务价值总结", font_name=TITLE_FONT, font_size=34, bold=True, color=WHITE)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.2), ACCENT_GOLD, Pt(3))

value_items = [
    ("80%", "情报效率提升", "从人工到AI自动聚合", ACCENT_BLUE),
    ("60-95%", "Token节省", "上下文压缩技术", GREEN_ACCENT),
    ("3-5分钟", "一键部署", "从模板到上线", ACCENT_GOLD),
    ("2", "个开源项目", "InsightPro + SAC", PURPLE_ACCENT),
    ("3", "项华为云服务", "MaaS·RFS·Headroom", ACCENT_CYAN),
    ("6+6+4", "多场景全覆盖", "行业+竞争+场景", RED_ACCENT),
]

for i, (num, title, desc, clr) in enumerate(value_items):
    x = 0.4 + (i % 3) * 4.3
    y = 1.4 + (i // 3) * 2.8
    # Card
    add_shape(slide, Inches(x), Inches(y), Inches(3.9), Inches(2.4),
              fill_color=RGBColor(0x16, 0x1B, 0x33), line_color=clr)
    # Number
    add_text_box(slide, Inches(x + 0.2), Inches(y + 0.2), Inches(3.5), Inches(0.8),
                 num, font_name=TITLE_FONT, font_size=38, bold=True, color=clr,
                 alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, Inches(x + 0.2), Inches(y + 1.0), Inches(3.5), Inches(0.4),
                 title, font_name=BODY_FONT, font_size=16, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    # Accent line
    add_accent_line(slide, Inches(x + 0.5), Inches(y + 1.5), Inches(2.9), clr, Pt(1))
    # Description
    add_text_box(slide, Inches(x + 0.2), Inches(y + 1.6), Inches(3.5), Inches(0.4),
                 desc, font_name=BODY_FONT, font_size=12, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Pt(4), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 27 - FUTURE PLANNING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(6), Inches(0.6),
             "未来规划", font_name=TITLE_FONT, font_size=34, bold=True, color=DARK_TEXT)
add_accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.2), ACCENT_GOLD, Pt(3))

# InsightPro future
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5),
          fill_color=DARK_BG, line_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(1.4), Inches(5), Inches(0.5),
             "InsightPro 规划", font_name=TITLE_FONT, font_size=20, bold=True, color=ACCENT_BLUE)
add_accent_line(slide, Inches(1.1), Inches(1.9), Inches(2), ACCENT_BLUE, Pt(2))

ip_future = [
    ("Supabase持久化", "替代SQLite，支持多用户", ACCENT_BLUE),
    ("用户认证系统", "RBAC权限管理", PURPLE_ACCENT),
    ("全局搜索", "全文检索能力", GREEN_ACCENT),
    ("通知系统", "站内信+邮件推送", ACCENT_CYAN),
    ("PDF导出", "报告生成与下载", ACCENT_GOLD),
    ("移动端适配", "响应式多端覆盖", RED_ACCENT),
]
for i, (f, desc, clr) in enumerate(ip_future):
    y = 2.2 + i * 0.7
    add_shape(slide, Inches(1.1), Inches(y), Pt(6), Pt(6), fill_color=clr, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(1.4), Inches(y - 0.05), Inches(1.5), Inches(0.3),
                 f, font_name=BODY_FONT, font_size=13, bold=True, color=WHITE)
    add_text_box(slide, Inches(3.0), Inches(y - 0.05), Inches(2.8), Inches(0.3),
                 desc, font_name=BODY_FONT, font_size=11, color=MID_GRAY)

# SAC future
add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5),
          fill_color=DARK_BG2, line_color=ACCENT_GOLD)
add_text_box(slide, Inches(7.3), Inches(1.4), Inches(5), Inches(0.5),
             "SAC 规划", font_name=TITLE_FONT, font_size=20, bold=True, color=ACCENT_GOLD)
add_accent_line(slide, Inches(7.3), Inches(1.9), Inches(2), ACCENT_GOLD, Pt(2))

sac_future = [
    ("更多场景覆盖", "持续扩展方案矩阵", ACCENT_BLUE),
    ("AI自动化增强", "智能编排与自愈", PURPLE_ACCENT),
    ("跨方案编排", "多方案协同部署", GREEN_ACCENT),
    ("质量评估体系", "自动化测试验证", ACCENT_CYAN),
    ("社区协作", "开源共建生态", ACCENT_GOLD),
    ("性能优化", "更极致的压缩比", RED_ACCENT),
]
for i, (f, desc, clr) in enumerate(sac_future):
    y = 2.2 + i * 0.7
    add_shape(slide, Inches(7.3), Inches(y), Pt(6), Pt(6), fill_color=clr, shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(7.6), Inches(y - 0.05), Inches(1.5), Inches(0.3),
                 f, font_name=BODY_FONT, font_size=13, bold=True, color=WHITE)
    add_text_box(slide, Inches(9.2), Inches(y - 0.05), Inches(2.8), Inches(0.3),
                 desc, font_name=BODY_FONT, font_size=11, color=MID_GRAY)

add_shape(slide, Inches(0), Inches(7.42), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

# ═══════════════════════════════════════════════════════════════
# SLIDE 28 - CLOSING PAGE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.15), fill_color=ACCENT_GOLD)

# Core quote
add_shape(slide, Inches(0), Inches(3.1), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)
add_shape(slide, Inches(0), Inches(4.3), Inches(13.33), Inches(0.08), fill_color=ACCENT_GOLD)

add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
             "AI Agent 不是替代人，而是放大人的能力",
             font_name=TITLE_FONT, font_size=32, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)

# GitHub links
add_text_box(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
             "GitHub", font_name=TITLE_FONT, font_size=16, bold=True, color=ACCENT_GOLD,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.4),
             "github.com/Justin-TangPan/Insightpro", font_name=BODY_FONT, font_size=14, bold=False, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.4),
             "github.com/Justin-TangPan/solution-practices", font_name=BODY_FONT, font_size=14, bold=False, color=LIGHT_GRAY,
             alignment=PP_ALIGN.CENTER)

# Speaker
add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.5),
             "唐潘 · 解决方案", font_name=BODY_FONT, font_size=18, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(6.6), Inches(10.3), Inches(0.4),
             "2026", font_name=TITLE_FONT, font_size=14, bold=False, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER)

# Decorative circles
for i, (x, y, sz) in enumerate([(1.0, 0.6, 0.4), (11.5, 0.8, 0.25), (1.2, 6.8, 0.2), (11.8, 6.5, 0.15)]):
    c = add_shape(slide, Inches(x), Inches(y), Inches(sz), Inches(sz),
                  fill_color=None, shape_type=MSO_SHAPE.OVAL)
    c.fill.solid()
    c.fill.fore_color.rgb = ACCENT_GOLD if i % 2 == 0 else ACCENT_BLUE
    c.line.fill.background()

add_shape(slide, Inches(0), Inches(7.35), Inches(13.33), Inches(0.15), fill_color=ACCENT_GOLD)

# ── SAVE ───────────────────────────────────────────────────────
output_dir = os.path.dirname(os.path.abspath(__file__))
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "competition.pptx")
prs.save(output_path)
print(f"[OK] PPT saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
