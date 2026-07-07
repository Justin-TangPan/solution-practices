#!/usr/bin/env python3
"""Add 2 detailed architecture pages to competition.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BG = RGBColor(0x0F, 0x0F, 0x23)
DARK_BG2 = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
PURPLE = RGBColor(0xA7, 0x8B, 0xF0)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
LAYER_BG = RGBColor(0xF1, 0xF5, 0xF9)
SUBTLE = RGBColor(0x94, 0xA3, 0xB8)

prs = Presentation(r'C:\Users\Administrator\Desktop\Project\claudeproject\solution-practices\ppt\competition\competition.pptx')

def set_bg(slide, color):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.width = bw
    else: s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.width = bw
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=Pt(12), color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = sz; p.font.color.rgb = color; p.font.bold = bold; p.font.name = font; p.alignment = align
    return box

def arrow_down(slide, x, y, w=Inches(0.3)):
    """Vertical down arrow using triangle shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, Inches(0.25))
    s.fill.solid(); s.fill.fore_color.rgb = SUBTLE; s.line.fill.background()
    return s

def arrow_right(slide, x, y, h=Inches(0.2)):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, Inches(0.25), h)
    s.fill.solid(); s.fill.fore_color.rgb = SUBTLE; s.line.fill.background()
    return s

def layer_block(slide, x, y, w, h, label, items, label_color=BLUE, item_color=WHITE, border_color=CARD_BORDER, is_dark=False):
    """A labeled group of items (like a layer in architecture)."""
    bg = DARK_BG2 if is_dark else LAYER_BG
    rrect(slide, x, y, w, h, fill=bg, border=border_color, bw=Pt(1))
    # Layer label
    rrect(slide, x + Inches(0.08), y + Inches(0.08), Inches(1.1), Inches(0.22), fill=label_color)
    tb(slide, x + Inches(0.12), y + Inches(0.08), Inches(1.0), Inches(0.22), label,
       Pt(7), WHITE, bold=True, font='Consolas')
    # Items
    ix = x + Inches(0.15)
    iy = y + Inches(0.38)
    iw = (w - Inches(0.3)) / len(items) - Inches(0.05) if len(items) > 0 else Inches(0.5)
    for i, item_text in enumerate(items):
        item_bg = WHITE if not is_dark else RGBColor(0x25, 0x25, 0x40)
        item_border = BLUE if is_dark else BLUE
        rrect(slide, ix + i * (iw + Inches(0.05)), iy, iw, h - Inches(0.45),
              fill=item_bg, border=item_border, bw=Pt(0.5))
        tb(slide, ix + i * (iw + Inches(0.05)) + Inches(0.05), iy + Inches(0.03),
           iw - Inches(0.1), h - Inches(0.5), item_text, Pt(7.5),
           WHITE if is_dark else DARK_TEXT)

# ==================================================================
# PAGE 29: InsightPro 技术架构（详细版）
# ==================================================================
s29 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s29, LIGHT_BG)

# Chrome
tb(s29, Inches(0.6), Inches(0.25), Inches(6), Inches(0.3), 'InsightPro · 详细技术架构', Pt(10), GRAY_TEXT, font='Consolas')
tb(s29, Inches(10), Inches(0.25), Inches(2.8), Inches(0.3), 'Technical Arch · 29 / 30', Pt(10), GRAY_TEXT, align=PP_ALIGN.RIGHT, font='Consolas')
rect(s29, Inches(0.6), Inches(0.6), Inches(12.1), Pt(1), fill=CARD_BORDER)

# Title
tb(s29, Inches(0.6), Inches(0.7), Inches(5), Inches(0.35), 'SYSTEM ARCHITECTURE', Pt(10), BLUE, bold=True, font='Consolas')
tb(s29, Inches(0.6), Inches(1.0), Inches(8), Inches(0.5), 'InsightPro 全栈技术架构', Pt(28), DARK_TEXT, bold=True)

# ---- Layer: 前端展示层 ----
layer_block(s29, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.0),
            'FRONTEND', [
                'Next.js 16.2\nApp Router\n15页面路由',
                'React 19.2\nTypeScript 5\nTailwind CSS 4',
                'Recharts 3.8\nAreaChart\nBarChart/Line',
                'ChatAssistant\n流式SSE对话\nRecharts图',
                'PageTracker\nsendBeacon\n路由埋点',
                'Sidebar+Mobile\n3 NavGroup\n响应式布局'
            ], label_color=BLUE, is_dark=False)

# Down arrow
arrow_down(s29, Inches(6.5), Inches(2.65))

# ---- Layer: API 网关层 ----
layer_block(s29, Inches(0.6), Inches(2.95), Inches(12.1), Inches(1.0),
            'API', [
                'FastAPI 0.136\nUvicorn 0.48\n31端点',
                'CORS中间件\nlocalhost\n94.74.90.21\n*.vercel.app',
                'SSE流式\n/api/chat/stream\ntext/event-',
                'HTTP REST\nGET/POST/DELETE\nJSON通信',
            ], label_color=GREEN, is_dark=False)

# Down arrow
arrow_down(s29, Inches(6.5), Inches(4.0))

# ---- Layer: 业务逻辑层 ----
layer_block(s29, Inches(0.6), Inches(4.35), Inches(12.1), Inches(1.15),
            'BUSINESS', [
                'GitHub爬虫\nBeautifulSoup\n每日09:00\ndaily/weekly',
                '百度热搜\n3级降级策略\n__INITIAL__\n5min轮询',
                'AI分析任务\nDeepSeek API\nOpenAI SDK\nJSON输出',
                '智能客服\nModelArts\ndeepseek-v3.2\n流式+非流',
                '招标/需求\n日期轮换机制\n预置数据池\n每日自动',
                '邮件服务\nQQ SMTP\n每日09:05\n订阅推送'
            ], label_color=GOLD, is_dark=False)

# Down arrow
arrow_down(s29, Inches(6.5), Inches(5.55))

# ---- Layer: RAG 引擎层 ----
rrect(s29, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.0), fill=RGBColor(0xEF,0xF6,0xFF), border=BLUE, bw=Pt(1))
rrect(s29, Inches(0.68), Inches(5.93), Inches(0.95), Inches(0.22), fill=BLUE)
tb(s29, Inches(0.72), Inches(5.93), Inches(0.9), Inches(0.22), 'RAG', Pt(7), WHITE, bold=True, font='Consolas')

# 5 vector collections
cols_rag = [
    ('Qdrant向量库\n本地文件持久化\ndeep_searcher_q', BLUE),
    ('bge-m3\nSentenceTransform\n1024维 中英双语', GREEN),
    ('5向量集合\nplatform_know\ndemand_signal\ncompetitive', PURPLE),
    ('github_trnd\npolicies\n跨集合检索\n按score排序', PURPLE),
    ('DeepSearcher\nzilliztech\nChainOfRAG\n迭代搜索', BLUE),
]
for i, (txt, c) in enumerate(cols_rag):
    x = Inches(0.68) + i * Inches(2.45)
    rrect(s29, x, Inches(6.25), Inches(2.25), Inches(0.52), fill=WHITE, border=c, bw=Pt(0.5))
    tb(s29, x + Inches(0.05), Inches(6.28), Inches(2.15), Inches(0.48), txt, Pt(7), DARK_TEXT)

# Down arrow
arrow_down(s29, Inches(6.5), Inches(6.9))

# ---- Layer: 数据存储层 ----
rrect(s29, Inches(0.6), Inches(7.2), Inches(12.1), Inches(0.0), fill=RGBColor(0xFF,0xF7,0xED), border=GOLD, bw=Pt(0))
# Using a shorter height
rrect(s29, Inches(0.6), Inches(7.2), Inches(12.1), Inches(0.65), fill=RGBColor(0xFF,0xF7,0xED), border=GOLD, bw=Pt(1))
rrect(s29, Inches(0.68), Inches(7.28), Inches(0.95), Inches(0.22), fill=GOLD)
tb(s29, Inches(0.72), Inches(7.28), Inches(0.9), Inches(0.22), 'DATA', Pt(7), WHITE, bold=True, font='Consolas')

db_items = [
    ('SQLite trending.db\n9张表 · 降级兜底', GOLD),
    ('Supabase PostgreSQL\nPrisma 7.8 · 预留扩展', BLUE),
    ('Qdrant文件存储\n5集合 · bge-m3索引', PURPLE),
    ('8张表详情\ngithub_trending\nbaidu_hotsearch\ncompetitor_news', GOLD),
    ('bidding_opps\ndemand_signals\npage_visits\nemail_subs', GOLD),
]
for i, (txt, c) in enumerate(db_items):
    x = Inches(0.68) + i * Inches(2.45)
    rrect(s29, x, Inches(7.58), Inches(2.25), Inches(0.22), fill=WHITE, border=c, bw=Pt(0.5))
    tb(s29, x + Inches(0.05), Inches(7.58), Inches(2.15), Inches(0.22), txt.split('\n')[0], Pt(7), DARK_TEXT)

# Right side: External services
rrect(s29, Inches(9.5), Inches(2.0), Inches(3.2), Inches(3.3), fill=WHITE, border=BLUE, bw=Pt(1))
tb(s29, Inches(9.6), Inches(2.05), Inches(3), Inches(0.25), '外部服务集成', Pt(9), BLUE, bold=True)
ext_items = [
    'DeepSeek API (api.deepseek.com)',
    '  -> 模型: deepseek-chat',
    '  -> OpenAI SDK 兼容调用',
    '  -> 业务分析/招标/市场',
    '',
    'ModelArts MaaS',
    '  -> api.modelarts-maas.com',
    '  -> 模型: deepseek-v3.2',
    '  -> httpx 直接REST调用',
    '  -> 流式SSE客服对话',
    '',
    'Supabase',
    '  -> PostgreSQL + SSR Auth',
    '  -> ap-southeast-1 区域',
    '  -> PgBouncer连接池 :6543',
    '',
    'QQ SMTP (smtp.qq.com:465)',
    '  -> 每日洞察邮件推送',
]
for i, item in enumerate(ext_items):
    clr = DARK_TEXT if not item.startswith(' ') else GRAY_TEXT
    bld = not item.startswith(' ') and item != ''
    tb(s29, Inches(9.6), Inches(2.35 + i * 0.175), Inches(3.0), Inches(0.17), item, Pt(7), clr, bold=bld)

# APScheduler info
rrect(s29, Inches(9.5), Inches(5.4), Inches(3.2), Inches(2.0), fill=WHITE, border=GOLD, bw=Pt(1))
tb(s29, Inches(9.6), Inches(5.45), Inches(3), Inches(0.25), 'APScheduler 定时任务', Pt(9), GOLD, bold=True)
cron_items = [
    '09:00  github_daily (爬GitHub)',
    '08:30  bidding_daily (招标采集)',
    '08:00  demand_daily (需求信号)',
    '09:05  daily_email (邮件推送)',
    '03:00  cleanup (清理30天数据)',
    '',
    '时区: Asia/Shanghai',
    '调度器: BackgroundScheduler',
]
for i, item in enumerate(cron_items):
    tb(s29, Inches(9.6), Inches(5.75 + i * 0.17), Inches(3.0), Inches(0.17), item, Pt(7), DARK_TEXT if item else GRAY_TEXT)

# Version tags
rrect(s29, Inches(0.6), Inches(7.92), Inches(12.1), Inches(0.22), fill=DARK_TEXT, border=DARK_TEXT)
ver_text = 'Python 3.13  |  FastAPI 0.136.3  |  Next.js 16.2.6  |  React 19.2.4  |  DeepSearcher  |  Qdrant  |  bge-m3  |  Prisma 7.8  |  SQLite'
tb(s29, Inches(0.8), Inches(7.93), Inches(11.5), Inches(0.2), ver_text, Pt(7.5), WHITE, font='Consolas', align=PP_ALIGN.CENTER)

# Footer
tb(s29, Inches(0.6), Inches(8.15), Inches(6), Inches(0.35), 'InsightPro · 详细技术架构', Pt(9), GRAY_TEXT, font='Consolas')
tb(s29, Inches(10), Inches(8.15), Inches(2.8), Inches(0.35), 'Technical Architecture', Pt(9), GRAY_TEXT, align=PP_ALIGN.RIGHT, font='Consolas')

print("  Slide 29: InsightPro Technical Architecture")


# ==================================================================
# PAGE 30: SAC 解决方案实践 架构（详细版）
# ==================================================================
s30 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s30, DARK_BG)

# Chrome
tb(s30, Inches(0.6), Inches(0.25), Inches(6), Inches(0.3), 'SAC · 解决方案实践架构', Pt(10), SUBTLE, font='Consolas')
tb(s30, Inches(10), Inches(0.25), Inches(2.8), Inches(0.3), 'Technical Arch · 30 / 30', Pt(10), SUBTLE, align=PP_ALIGN.RIGHT, font='Consolas')
rect(s30, Inches(0.6), Inches(0.6), Inches(12.1), Pt(1), fill=RGBColor(0x33,0x33,0x50))

# Title
tb(s30, Inches(0.6), Inches(0.7), Inches(5), Inches(0.35), 'SOLUTION PRACTICES ARCHITECTURE', Pt(10), GOLD, bold=True, font='Consolas')
tb(s30, Inches(0.6), Inches(1.0), Inches(8), Inches(0.5), 'SAC 解决方案实践 · 全栈技术架构', Pt(28), WHITE, bold=True)

# ---- 7 Solutions Layer ----
rrect(s30, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.85), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60), bw=Pt(1))
rrect(s30, Inches(0.68), Inches(1.68), Inches(1.1), Inches(0.22), fill=GOLD)
tb(s30, Inches(0.72), Inches(1.68), Inches(1.0), Inches(0.22), 'SOLUTIONS', Pt(7), WHITE, bold=True, font='Consolas')

sols = [
    ('Headroom\nOpenCode', 'HK pip', BLUE),
    ('Headroom\nClaude Code', 'HK pip', BLUE),
    ('LiteLLM', 'CN/HK Docker', GREEN),
    ('Supabase', 'CN Docker', GREEN),
    ('OpenHands', 'CN Docker', PURPLE),
    ('CodeWhale', 'CN Docker', PURPLE),
    ('AiToEarn', 'CN/HK Docker', GOLD),
]
for i, (name, mode, c) in enumerate(sols):
    x = Inches(0.68) + i * Inches(1.73)
    rrect(s30, x, Inches(1.98), Inches(1.6), Inches(0.4), fill=RGBColor(0x25,0x25,0x40), border=c, bw=Pt(0.5))
    tb(s30, x + Inches(0.05), Inches(1.99), Inches(1.5), Inches(0.22), name, Pt(7), WHITE, bold=True)
    tb(s30, x + Inches(0.05), Inches(2.2), Inches(1.5), Inches(0.15), mode, Pt(6), SUBTLE)

# Down arrow
arrow_down(s30, Inches(6.5), Inches(2.5))

# ---- IaC 模板层 ----
rrect(s30, Inches(0.6), Inches(2.85), Inches(12.1), Inches(1.1), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60), bw=Pt(1))
rrect(s30, Inches(0.68), Inches(2.93), Inches(1.5), Inches(0.22), fill=GREEN)
tb(s30, Inches(0.72), Inches(2.93), Inches(1.4), Inches(0.22), 'TERRAFORM', Pt(7), WHITE, bold=True, font='Consolas')

tf_items = [
    ('Provider\nhuaweicloud\nregion only\n(Pitfall#19)', GREEN),
    ('Variables\n7个核心变量\nsolution_name\necs_flavor', BLUE),
    ('Locals\nname_suffix\nsubstr(uuid(),0,8)\n唯一命名', BLUE),
    ('Resources\nVPC/Subnet\nSecGroup/EIP\nECS Instance', PURPLE),
    ('user_data\n重置密码\nOBS下载脚本\n执行→清理', GOLD),
    ('.extension\n中英双语\n参数分组\nRFS面板', SUBTLE),
]
for i, (txt, c) in enumerate(tf_items):
    x = Inches(0.68) + i * Inches(2.03)
    rrect(s30, x, Inches(3.22), Inches(1.9), Inches(0.65), fill=RGBColor(0x25,0x25,0x40), border=c, bw=Pt(0.5))
    tb(s30, x + Inches(0.05), Inches(3.25), Inches(1.8), Inches(0.6), txt, Pt(7), WHITE)

# Down arrow
arrow_down(s30, Inches(6.5), Inches(4.0))

# ---- 安装脚本层 ----
rrect(s30, Inches(0.6), Inches(4.35), Inches(12.1), Inches(1.05), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60), bw=Pt(1))
rrect(s30, Inches(0.68), Inches(4.43), Inches(1.3), Inches(0.22), fill=GOLD)
tb(s30, Inches(0.72), Inches(4.43), Inches(1.2), Inches(0.22), 'INSTALL.SH', Pt(7), WHITE, bold=True, font='Consolas')

sh_items = [
    ('Stage 1\n系统准备\ndpkg configure\napt update', GOLD),
    ('Stage 2\nDocker安装\nCN:华为云镜像\nHK:官方源', BLUE),
    ('Stage 3\n应用配置\n目录/权限/配置\n备份+crontab', GREEN),
    ('Stage 4\n启动验证\ndocker up -d\n健康检查120s', PURPLE),
    ('Stage 5\nOpenCode配置\nopencode.json\n环境变量注入', BLUE),
]
for i, (txt, c) in enumerate(sh_items):
    x = Inches(0.68) + i * Inches(2.43)
    rrect(s30, x, Inches(4.72), Inches(2.3), Inches(0.6), fill=RGBColor(0x25,0x25,0x40), border=c, bw=Pt(0.5))
    tb(s30, x + Inches(0.05), Inches(4.75), Inches(2.2), Inches(0.55), txt, Pt(7), WHITE)

# Down arrow
arrow_down(s30, Inches(6.5), Inches(5.45))

# ---- 华为云基础设施层 ----
rrect(s30, Inches(0.6), Inches(5.75), Inches(12.1), Inches(0.85), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60), bw=Pt(1))
rrect(s30, Inches(0.68), Inches(5.83), Inches(0.95), Inches(0.22), fill=BLUE)
tb(s30, Inches(0.72), Inches(5.83), Inches(0.9), Inches(0.22), 'ECS', Pt(7), WHITE, bold=True, font='Consolas')

infra = [
    ('VPC\n172.16.0.0/16', BLUE),
    ('Subnet\n172.16.1.0/24', BLUE),
    ('安全组\nSSH:22\nProxy:8787\nPing:ICMP', GREEN),
    ('EIP弹性IP\n300Mbit/s\n按流量计费', GOLD),
    ('Flexus X ECS\nx1.2u.4g\nUbuntu 24.04\n40GB SSD', PURPLE),
    ('OBS存储\n模板+脚本托管\nCN/HK双桶', GOLD),
]
for i, (txt, c) in enumerate(infra):
    x = Inches(0.68) + i * Inches(2.03)
    rrect(s30, x, Inches(6.12), Inches(1.9), Inches(0.4), fill=RGBColor(0x25,0x25,0x40), border=c, bw=Pt(0.5))
    tb(s30, x + Inches(0.05), Inches(6.14), Inches(1.8), Inches(0.36), txt, Pt(7), WHITE)

# Down arrow
arrow_down(s30, Inches(6.5), Inches(6.6))

# ---- 运行时层 ----
rrect(s30, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.65), fill=DARK_BG2, border=RGBColor(0x40,0x40,0x60), bw=Pt(1))
rrect(s30, Inches(0.68), Inches(6.98), Inches(1.0), Inches(0.22), fill=PURPLE)
tb(s30, Inches(0.72), Inches(6.98), Inches(0.95), Inches(0.22), 'RUNTIME', Pt(7), WHITE, bold=True, font='Consolas')

runtime = [
    ('OpenCode CLI\nnpm install -g\nopencode-ai@latest', BLUE),
    ('Headroom Proxy\npip install headroom-ai\nport 8787', GREEN),
    ('Context Compress\nAST感知\n60-95%压缩', PURPLE),
    ('MaaS API\napi.modelarts-maas.com\nDeepSeek/Claude', GOLD),
    ('Docker Compose\nLiteLLM+PG+Prom\nSupabase+Kong', BLUE),
]
for i, (txt, c) in enumerate(runtime):
    x = Inches(0.68) + i * Inches(2.43)
    rrect(s30, x, Inches(7.25), Inches(2.3), Inches(0.25), fill=RGBColor(0x25,0x25,0x40), border=c, bw=Pt(0.5))
    tb(s30, x + Inches(0.05), Inches(7.26), Inches(2.2), Inches(0.22), txt.split('\n')[0], Pt(7), WHITE)

# Right side: AI Skills
rrect(s30, Inches(9.5), Inches(2.0), Inches(3.2), Inches(2.0), fill=RGBColor(0x25,0x25,0x40), border=GOLD, bw=Pt(1))
tb(s30, Inches(9.6), Inches(2.05), Inches(3), Inches(0.25), 'AI Skills 技能体系', Pt(9), GOLD, bold=True)
skill_items = [
    'rfs-practices',
    '  创建/修复RFS模板',
    '  10决策点+27坑',
    '',
    'solution-extractor',
    '  爬取方案页→Excel',
    '  agent-browser链',
    '',
    'page-enhance',
    '  方案页面增强',
    '  10维度优化检查',
    '',
    'deep-search-synth',
    '  3层深度搜索',
    '  多源交叉验证',
]
for i, item in enumerate(skill_items):
    clr = WHITE if not item.startswith(' ') else SUBTLE
    bld = not item.startswith(' ') and item != ''
    tb(s30, Inches(9.6), Inches(2.35 + i * 0.155), Inches(3.0), Inches(0.15), item, Pt(7), clr, bold=bld)

# Right side: Scripts
rrect(s30, Inches(9.5), Inches(4.1), Inches(3.2), Inches(1.6), fill=RGBColor(0x25,0x25,0x40), border=GREEN, bw=Pt(1))
tb(s30, Inches(9.6), Inches(4.15), Inches(3), Inches(0.25), 'Python 工具链', Pt(9), GREEN, bold=True)
script_items = [
    'validate_template.py',
    '  模板结构验证器',
    '',
    'gen_docx.py',
    '  README→Word文档',
    '',
    'gen_sac_docx.py',
    '  SAC报告生成',
    '',
    'generate_extension.py',
    '  .extension配置生成',
    '',
    'package_solution.sh',
    '  方案打包归档',
]
for i, item in enumerate(script_items):
    clr = WHITE if not item.startswith(' ') else SUBTLE
    bld = not item.startswith(' ') and item != ''
    tb(s30, Inches(9.6), Inches(4.45 + i * 0.115), Inches(3.0), Inches(0.15), item, Pt(6.5), clr, bold=bld)

# Right side: Regions
rrect(s30, Inches(9.5), Inches(5.8), Inches(3.2), Inches(1.3), fill=RGBColor(0x25,0x25,0x40), border=BLUE, bw=Pt(1))
tb(s30, Inches(9.6), Inches(5.85), Inches(3), Inches(0.25), '区域部署差异', Pt(9), BLUE, bold=True)
region_items = [
    'CN (国内)',
    '  Docker源:华为云镜像',
    '  SWR:需预推镜像',
    '  PyPI:清华源',
    '  部署:OBS下载脚本',
    '',
    'HK (海外)',
    '  Docker源:官方',
    '  SWR:不需要',
    '  PyPI:默认源',
    '  部署:内联user_data',
]
for i, item in enumerate(region_items):
    clr = GOLD if not item.startswith(' ') else SUBTLE
    bld = not item.startswith(' ') and item != ''
    tb(s30, Inches(9.6), Inches(6.15 + i * 0.115), Inches(3.0), Inches(0.15), item, Pt(6.5), clr, bold=bld)

# Version bar
rrect(s30, Inches(0.6), Inches(7.62), Inches(12.1), Inches(0.22), fill=RGBColor(0x25,0x25,0x40), border=RGBColor(0x40,0x40,0x60))
ver_sac = 'Terraform HCL/JSON  |  Docker Compose  |  Flexus X ECS  |  VPC+SecGroup+EIP  |  OBS  |  Python 3 + Bash  |  4 Skills'
tb(s30, Inches(0.8), Inches(7.63), Inches(11.5), Inches(0.2), ver_sac, Pt(7.5), SUBTLE, font='Consolas', align=PP_ALIGN.CENTER)

# Footer
tb(s30, Inches(0.6), Inches(7.9), Inches(6), Inches(0.35), 'SAC · 详细技术架构', Pt(9), SUBTLE, font='Consolas')
tb(s30, Inches(10), Inches(7.9), Inches(2.8), Inches(0.35), 'Solution Practices', Pt(9), SUBTLE, align=PP_ALIGN.RIGHT, font='Consolas')

print("  Slide 30: SAC Technical Architecture")

# ---- SAVE ----
out = r'C:\Users\Administrator\Desktop\Project\claudeproject\solution-practices\ppt\competition\competition.pptx'
prs.save(out)
import os
print(f"\nDone! Total slides: {len(prs.slides)}, Size: {os.path.getsize(out)/1024:.0f} KB")
